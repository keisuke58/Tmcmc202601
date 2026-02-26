#!/usr/bin/env python3
"""
論文（Markdown）を PowerPoint (.pptx) に変換するスクリプト.

paper_summary_ja.md などの Markdown ファイルを読み込み、
セクション・箇条書き・表・画像をスライドに変換する。

Usage:
    python paper_to_pptx.py [input.md] [output.pptx]
    python paper_to_pptx.py  # デフォルト: paper_summary_ja.md → paper_slides.pptx
"""

import re
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# デザイン定数（アカデミック・プロフェッショナル）
COLORS = {
    "primary": RGBColor(0x1E, 0x3A, 0x5F),  # ネイビー
    "accent": RGBColor(0x0D, 0x73, 0x77),  # ティール
    "section_bg": RGBColor(0x2C, 0x3E, 0x50),  # ダークスレート
    "text": RGBColor(0x2C, 0x3E, 0x50),
    "text_light": RGBColor(0x7F, 0x8C, 0x8D),
}

# デフォルトパス
DEFAULT_INPUT = Path(__file__).resolve().parent.parent / "docs" / "paper_summary_ja.md"
DEFAULT_OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "paper_slides.pptx"


def _parse_markdown(content: str, md_path: Path) -> list[dict]:
    """
    Markdown をパースしてスライド用の構造に変換する.

    Parameters
    ----------
    content : str
        Markdown 本文
    md_path : Path
        Markdown ファイルのパス（画像相対パス解決用）

    Returns
    -------
    list[dict]
        各要素: {"type", "text", "items", "tables", "images"}
    """
    slides_data: list[dict] = []
    current_slide: dict = {}
    current_items: list[str] = []
    in_table = False
    table_rows: list[list[str]] = []

    lines = content.split("\n")

    for line in lines:
        stripped = line.strip()

        # 区切り線・空行・数式ブロックはスキップ
        if not stripped or stripped == "---" or stripped.startswith("$$"):
            continue

        # 画像: ![alt](path)
        img_match = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)$", stripped)
        if img_match:
            alt, path = img_match.group(1), img_match.group(2)
            resolved = (md_path.parent / path).resolve()
            current_slide.setdefault("images", []).append(
                {
                    "path": resolved,
                    "alt": alt or "Figure",
                }
            )
            continue

        # 表の処理
        if re.match(r"^\|.+\|$", stripped):
            if not in_table:
                in_table = True
                table_rows = []
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if cells and not all(re.match(r"^[-:]+$", c) for c in cells):
                table_rows.append(cells)
            continue
        else:
            if in_table and table_rows:
                if current_slide:
                    current_slide.setdefault("tables", []).append(table_rows)
                else:
                    slides_data.append({"type": "content", "tables": [table_rows]})
                table_rows = []
            in_table = False

        # 見出し
        h_match = re.match(r"^(#{1,3})\s+(.+)$", stripped)
        if h_match:
            level = len(h_match.group(1))
            title = _clean_markdown(h_match.group(2).strip())

            if current_slide:
                if current_items:
                    current_slide["items"] = current_items
                slides_data.append(current_slide)
                current_items = []

            if level == 1:
                current_slide = {"type": "title", "text": title, "items": []}
            else:
                current_slide = {"type": "section", "level": level, "text": title, "items": []}
            continue

        # 箇条書き
        if stripped.startswith("- ") or stripped.startswith("* "):
            item = _clean_markdown(stripped[2:].strip())
            current_items.append(item)
            continue

        # 番号付きリスト
        num_match = re.match(r"^(\d+)\.\s+(.+)$", stripped)
        if num_match:
            current_items.append(_clean_markdown(num_match.group(2).strip()))
            continue

        # 通常テキスト
        if not stripped.startswith("```"):
            if stripped.startswith(">"):
                stripped = stripped[1:].strip()
            current_items.append(_clean_markdown(stripped))

    if current_slide:
        if current_items:
            current_slide["items"] = current_items
        slides_data.append(current_slide)

    return slides_data


def _clean_markdown(text: str) -> str:
    """Markdown 記法を除去."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


def _add_table(slide, rows: list[list[str]], left: float, top: float, width: float) -> float:
    """PowerPoint の表を追加。戻り値は次の top 位置."""
    if not rows:
        return top
    nrows, ncols = len(rows), max(len(r) for r in rows)
    row_h = Inches(0.35)
    col_w = width / ncols
    table = slide.shapes.add_table(nrows, ncols, left, top, width, nrows * row_h).table
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            if j < ncols:
                c = table.cell(i, j)
                c.text = cell[:80] if len(cell) > 80 else cell
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.name = "Meiryo"
                if i == 0:
                    for p in c.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    c.fill.solid()
                    c.fill.fore_color.rgb = COLORS["primary"]
    return top + nrows * row_h + Inches(0.2)


def _add_slide(
    prs: Presentation,
    slide_data: dict,
    md_path: Path,
) -> None:
    """1スライドを追加."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    left = Inches(0.6)
    top = Inches(0.5)
    width = Inches(8.8)
    content_width = Inches(8.8)
    max_height = Inches(6.5)

    # タイトルスライド
    if slide_data["type"] == "title":
        # メインタイトル
        tf = slide.shapes.add_textbox(left, Inches(2), width, Inches(1.2))
        p = tf.text_frame.paragraphs[0]
        p.text = slide_data["text"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER
        tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        items = slide_data.get("items", [])
        if items:
            sub_tf = slide.shapes.add_textbox(left, Inches(3.3), width, Inches(1.5))
            sub_tf.text_frame.word_wrap = True
            for item in items[:5]:  # サブタイトルは最大5行
                p = sub_tf.text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS["text_light"]
                p.font.name = "Meiryo"
                p.alignment = PP_ALIGN.CENTER
        return

    # セクション区切りスライド（見出しのみ、項目なし）
    if (
        slide_data["type"] == "section"
        and not slide_data.get("items")
        and not slide_data.get("tables")
        and not slide_data.get("images")
    ):
        # フル幅のセクションバー
        shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0),
            Inches(2.5),
            Inches(10),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["section_bg"]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = slide_data["text"]
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].font.name = "Meiryo"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return

    # セクション見出し（content タイプはタイトルなし）
    if slide_data.get("text"):
        tf = slide.shapes.add_textbox(left, top, width, Inches(0.7))
        p = tf.text_frame.paragraphs[0]
        p.text = slide_data["text"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS["accent"]
        p.font.name = "Meiryo"
        top += Inches(0.8)

    # 画像（最初の画像を大きく表示）
    images = slide_data.get("images", [])
    if images:
        img_path = images[0]["path"]
        if img_path.exists():
            try:
                pic = slide.shapes.add_picture(
                    str(img_path),
                    left,
                    top,
                    width=Inches(5.5),
                )
                top += Inches(pic.height / 914400) + Inches(0.2)
            except Exception as e:
                logger.warning("Image load failed %s: %s", img_path, e)
        # 残りの画像はスキップ（1スライド1図を推奨）

    # 表（PowerPoint ネイティブ表）
    for table_rows in slide_data.get("tables", []):
        top = _add_table(slide, table_rows, left, top, content_width)

    # 本文
    items = slide_data.get("items", [])
    if items:
        tf = slide.shapes.add_textbox(left, top, width, max_height - top - Inches(0.2))
        text_frame = tf.text_frame
        text_frame.word_wrap = True

        for item in items:
            if len(item) > 250:
                for chunk in [item[i : i + 250] for i in range(0, len(item), 250)]:
                    p = text_frame.add_paragraph()
                    p.text = chunk
                    p.font.size = Pt(13)
                    p.font.name = "Meiryo"
                    p.font.color.rgb = COLORS["text"]
                    p.space_after = Pt(4)
            else:
                p = text_frame.add_paragraph()
                p.text = "• " + item if not item.startswith("•") else item
                p.font.size = Pt(13)
                p.font.name = "Meiryo"
                p.font.color.rgb = COLORS["text"]
                p.space_after = Pt(4)
                p.level = 0


def _merge_section_slides(section_slides: list[dict]) -> list[dict]:
    """
    セクション区切りスライド（見出しのみ）を次のスライドにマージする.
    同じ ## レベルの連続する ### サブセクションも結合する.
    """
    if not section_slides:
        return []

    merged: list[dict] = []
    i = 0

    while i < len(section_slides):
        curr = section_slides[i]

        # セクション区切り（見出しのみ・項目なし）→ 次のスライドにマージ
        if (
            curr["type"] == "section"
            and not curr.get("items")
            and not curr.get("tables")
            and not curr.get("images")
            and i + 1 < len(section_slides)
        ):
            next_slide = section_slides[i + 1]
            merged.append(
                {
                    "type": "section",
                    "text": curr["text"],
                    "level": curr.get("level", 2),
                    "items": next_slide.get("items", []),
                    "tables": next_slide.get("tables", []),
                    "images": next_slide.get("images", []),
                }
            )
            i += 2
            continue

        # 同じ ## の連続 ### で、両方とも小さい場合は結合
        if (
            i + 1 < len(section_slides)
            and curr.get("level") == 3
            and section_slides[i + 1].get("level") == 3
            and len(curr.get("items", [])) <= 4
            and len(section_slides[i + 1].get("items", [])) <= 4
            and not curr.get("images")
            and not section_slides[i + 1].get("images")
        ):
            next_slide = section_slides[i + 1]
            merged.append(
                {
                    "type": "section",
                    "text": curr["text"] + " / " + next_slide.get("text", ""),
                    "level": 3,
                    "items": curr.get("items", []) + next_slide.get("items", []),
                    "tables": curr.get("tables", []) + next_slide.get("tables", []),
                    "images": [],
                }
            )
            i += 2
            continue

        merged.append(curr)
        i += 1

    return merged


def convert_md_to_pptx(
    input_path: Path,
    output_path: Path,
) -> None:
    """
    Markdown を PowerPoint に変換する.

    Parameters
    ----------
    input_path : Path
        入力 Markdown ファイル
    output_path : Path
        出力 .pptx ファイル
    """
    content = input_path.read_text(encoding="utf-8")
    slides_data = _parse_markdown(content, input_path)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slides = [s for s in slides_data if s["type"] == "title"]
    section_slides = [s for s in slides_data if s["type"] in ("section", "content")]

    # 最初のタイトルに「論文タイトル」の内容をマージ（サブタイトルとして）
    if title_slides and section_slides and section_slides[0].get("text") == "論文タイトル":
        title_slides[0].setdefault("items", []).extend(section_slides[0].get("items", []))
        section_slides = section_slides[1:]  # 論文タイトルスライドはスキップ

    for slide_data in title_slides:
        _add_slide(prs, slide_data, input_path)

    # スライドマージ: セクション区切り + 次のコンテンツを結合、chunk を大きく
    merged_slides = _merge_section_slides(section_slides)
    chunk_size = 12  # 1スライドあたりの項目数を増やして枚数削減

    for slide_data in merged_slides:
        items = slide_data.get("items", [])
        tables = slide_data.get("tables", [])
        images = slide_data.get("images", [])

        # 画像があるスライドは分割しない（1図1スライド）
        if images and len(images) >= 1:
            _add_slide(prs, slide_data, input_path)
        elif items and len(items) > chunk_size:
            for i in range(0, len(items), chunk_size):
                chunk = items[i : i + chunk_size]
                sub_slide = {**slide_data, "items": chunk}
                if i == 0:
                    sub_slide["tables"] = slide_data.get("tables", [])
                    sub_slide["images"] = slide_data.get("images", [])
                else:
                    sub_slide["tables"] = []
                    sub_slide["images"] = []
                _add_slide(prs, sub_slide, input_path)
        else:
            _add_slide(prs, slide_data, input_path)

    prs.save(str(output_path))
    logger.info("Saved: %s (%d slides)", output_path, len(prs.slides))


def main() -> None:
    """CLI エントリポイント."""
    import argparse

    parser = argparse.ArgumentParser(description="論文 Markdown を PowerPoint に変換")
    parser.add_argument("input", nargs="?", default=str(DEFAULT_INPUT), help="入力 Markdown")
    parser.add_argument("output", nargs="?", default=str(DEFAULT_OUTPUT), help="出力 .pptx")
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO)

    input_path = Path(args.input)
    output_path = Path(args.output)

    if not input_path.exists():
        logger.error("Input file not found: %s", input_path)
        raise SystemExit(1)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    convert_md_to_pptx(input_path, output_path)
    print(f"Done: {output_path}")


if __name__ == "__main__":
    main()
