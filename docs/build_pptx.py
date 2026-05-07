#!/usr/bin/env python3
"""Convert slides_nishioka_heine2025.pdf to PPTX.

Strategy: render each PDF page as a high-resolution PNG (300 DPI via PyMuPDF)
and embed it as a full-slide image in python-pptx. This preserves 100% of the
LaTeX/Beamer visual quality — math, TikZ, fonts — in PowerPoint format.
"""
import io
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

DOCS = Path(__file__).parent
PDF = DOCS / "slides_nishioka_heine2025.pdf"
OUT = DOCS / "slides_nishioka_heine2025.pptx"

DPI = 300  # render resolution (300 = print quality)
W_IN = 13.33  # 16:9 slide width in inches
H_IN = 7.50  # 16:9 slide height in inches

# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)

blank_layout = prs.slide_layouts[6]  # completely blank layout

doc = fitz.open(str(PDF))
print(f"PDF: {PDF.name}  ({len(doc)} pages)")

for page_num, page in enumerate(doc):
    mat = fitz.Matrix(DPI / 72, DPI / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)

    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=False)
    buf.seek(0)

    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(W_IN), Inches(H_IN))

    if (page_num + 1) % 10 == 0:
        print(f"  page {page_num + 1}/{len(doc)}")

doc.close()
prs.save(str(OUT))
print(f"Saved: {OUT}  ({len(prs.slides)} slides, {OUT.stat().st_size // 1024 // 1024} MB)")
