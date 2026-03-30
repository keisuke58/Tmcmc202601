#!/usr/bin/env python3
"""Build slides_nishioka.pptx from paper content using python-pptx.
Figures are rendered from PDF/PNG via PyMuPDF and embedded as images.
"""
import io
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Paths ───────────────────────────────────────────────────────────────────
DOCS = Path(__file__).parent
FIGS = DOCS / "figures"
OUT = DOCS / "slides_nishioka.pptx"

# ─── Constants ───────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)  # 16:9
IKM = RGBColor(0, 82, 147)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 30, 30)
LIGHT_BG = RGBColor(230, 240, 250)
MID_GRAY = RGBColor(90, 90, 90)
FONT = "Calibri"
HEADER_H = Inches(0.68)
MARGIN_X = Inches(0.42)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def pdf_to_png_bytes(pdf_path: Path, dpi: int = 260) -> bytes:
    """Render first page of a PDF to PNG bytes."""
    doc = fitz.open(str(pdf_path))
    page = doc[0]
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    doc.close()
    return buf.read()


def img_bytes(path: Path, dpi: int = 260) -> bytes:
    if path.suffix.lower() == ".pdf":
        return pdf_to_png_bytes(path, dpi)
    buf = io.BytesIO()
    Image.open(path).convert("RGB").save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


def add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def header(slide, text, size=28):
    """Blue header bar at top."""
    tf_box = slide.shapes.add_textbox(Inches(0), Inches(0), W, HEADER_H)
    tf_box.fill.solid()
    tf_box.fill.fore_color.rgb = IKM
    tf = tf_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(size)
    run.font.name = FONT
    run.font.color.rgb = WHITE


def body_box(slide, text, left, top, width, height, size=16, bold=False, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def bullet_box(slide, items, left, top, width, height, size=15, title=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(size + 1)
        run.font.name = FONT
        run.font.color.rgb = IKM
        tf.add_paragraph()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (i == 0 and not title) else tf.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(size)
        run.font.name = FONT
        run.font.color.rgb = DARK
    return tf


def add_figure(slide, img_data, left, top, width, height=None):
    buf = io.BytesIO(img_data)
    if height:
        slide.shapes.add_picture(buf, left, top, width, height)
    else:
        slide.shapes.add_picture(buf, left, top, width)


def bg_rect(slide, left, top, width, height, color=LIGHT_BG):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = IKM


def add_network_diagram(slide, left, top, width, height):
    nodes = {
        "So": {"x": 0.50, "y": 0.16, "fill": RGBColor(210, 230, 255)},
        "An": {"x": 0.18, "y": 0.46, "fill": RGBColor(210, 245, 210)},
        "Vei": {"x": 0.50, "y": 0.46, "fill": RGBColor(255, 250, 205)},
        "Fn": {"x": 0.82, "y": 0.46, "fill": RGBColor(235, 220, 250)},
        "Pg": {"x": 0.50, "y": 0.78, "fill": RGBColor(255, 210, 210)},
    }

    node_size = min(width, height) * 0.20
    shapes = {}
    for label, spec in nodes.items():
        cx = left + width * spec["x"]
        cy = top + height * spec["y"]
        shp = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx - node_size / 2,
            cy - node_size / 2,
            node_size,
            node_size,
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = spec["fill"]
        shp.line.color.rgb = MID_GRAY
        shp.line.width = Pt(1.5)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = DARK
        shapes[label] = shp

    def connect(a, b, text, color, dashed=False):
        ax = shapes[a].left + shapes[a].width / 2
        ay = shapes[a].top + shapes[a].height / 2
        bx = shapes[b].left + shapes[b].width / 2
        by = shapes[b].top + shapes[b].height / 2
        con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
        con.line.color.rgb = color
        con.line.width = Pt(2 if not dashed else 1.25)
        if dashed:
            con.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

        tx = min(ax, bx) + abs(ax - bx) / 2
        ty = min(ay, by) + abs(ay - by) / 2
        tb = slide.shapes.add_textbox(tx - Inches(0.6), ty - Inches(0.15), Inches(1.2), Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(12)
        r.font.color.rgb = MID_GRAY

    connect("So", "An", "co-agg", RGBColor(0, 90, 160))
    connect("So", "Vei", "lactate", RGBColor(0, 90, 160))
    connect("So", "Fn", "formate", RGBColor(0, 90, 160))
    connect("Vei", "Pg", "pH", RGBColor(0, 90, 160))
    connect("Fn", "Pg", "peptides", RGBColor(0, 90, 160))
    connect("An", "Vei", "", RGBColor(150, 150, 150), dashed=True)
    connect("An", "Fn", "", RGBColor(150, 150, 150), dashed=True)
    connect("Vei", "Fn", "", RGBColor(150, 150, 150), dashed=True)
    connect("So", "Pg", "", RGBColor(150, 150, 150), dashed=True)
    connect("An", "Pg", "", RGBColor(150, 150, 150), dashed=True)


# ─── Build presentation ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Slide 1: Title ─────────────────────────────────────────────────────────
sl = add_slide(prs)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(240, 245, 252)
bg.line.fill.background()
bar = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.12))
bar.fill.solid()
bar.fill.fore_color.rgb = IKM
bar.line.fill.background()
bar2 = sl.shapes.add_shape(1, Inches(0), Inches(7.38), W, Inches(0.12))
bar2.fill.solid()
bar2.fill.fore_color.rgb = IKM
bar2.line.fill.background()

tb = sl.shapes.add_textbox(MARGIN_X, Inches(1.35), W - 2 * MARGIN_X, Inches(2.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = (
    "GPU-Accelerated Bayesian Inference of\nMulti-Species Biofilm Interaction Parameters via TMCMC"
)
r.font.size = Pt(34)
r.font.bold = True
r.font.name = FONT
r.font.color.rgb = IKM

tb2 = sl.shapes.add_textbox(MARGIN_X, Inches(3.85), W - 2 * MARGIN_X, Inches(0.6))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "15-Dimensional A-Matrix Estimation (µᵢ-free model)"
r2.font.size = Pt(18)
r2.font.name = FONT
r2.font.color.rgb = DARK

tb3 = sl.shapes.add_textbox(MARGIN_X, Inches(4.75), W - 2 * MARGIN_X, Inches(1.8))
tf3 = tb3.text_frame
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = (
    "Keisuke Nishioka,  Felix Klempt,  Hendrik Geisler,  "
    "Meisam Soleimani,  Philipp Junker\n"
    "Institute of Continuum Mechanics, Leibniz Universität Hannover · 2026"
)
r3.font.size = Pt(16)
r3.font.name = FONT
r3.font.color.rgb = RGBColor(60, 60, 60)

# ── Slide 2: Motivation ───────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Motivation & Background: Peri-Implantitis")

bullet_box(
    sl,
    [
        "Peri-implantitis: inflammatory destruction of implant tissue",
        "Cause: multi-species oral biofilm dysbiosis",
        "5 key species: So, An, Vei, Fn, Pg",
        "Heine et al. (2025): 4 conditions × 6 timepoints, N=8 replicates",
        "Commensal vs Dysbiotic  ×  Static vs HOBIC reactor",
        "Key: avirulent Pg ATCC 20709 (commensal) vs virulent Pg W83 (dysbiotic)",
    ],
    Inches(0.3),
    Inches(0.75),
    Inches(5.5),
    Inches(5.5),
    size=16,
    title="Clinical Problem",
)

bullet_box(
    sl,
    [
        "15-dim symmetric interaction matrix A  (5×5, Aᵢⱼ = Aⱼᵢ)",
        "GPU-accelerated TMCMC on RTX 4090",
        "Two-phase strategy: fix-ψ Phase 1 → free-ψ Phase 2",
        "Full discovery: all 15 params estimated without a priori locking",
    ],
    Inches(0.3),
    Inches(4.2),
    Inches(5.5),
    Inches(3.0),
    size=16,
    title="This Work",
)

bg_rect(sl, Inches(5.9), Inches(0.75), Inches(7.1), Inches(6.5))
body_box(
    sl,
    "Species Interaction Network",
    Inches(6.0),
    Inches(0.85),
    Inches(6.8),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)
add_network_diagram(sl, Inches(6.2), Inches(1.45), Inches(6.6), Inches(4.5))
body_box(
    sl,
    "Solid: confirmed pathways   |   Dashed: estimated freely",
    Inches(6.2),
    Inches(6.10),
    Inches(6.6),
    Inches(0.35),
    size=12,
    color=MID_GRAY,
)

# ── Slide 3: Mathematical Model ───────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Mathematical Model: Extended Hamilton Principle")

body_box(
    sl,
    "State variables (per material point):",
    Inches(0.3),
    Inches(0.8),
    Inches(6.0),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)
bullet_box(
    sl,
    [
        "φᵢ ∈ [0,1] — volume fraction of species i",
        "ψᵢ ∈ [0,1] — viability fraction (membrane-intact)",
        "φ̄ᵢ = φᵢψᵢ   — living volume fraction",
    ],
    Inches(0.3),
    Inches(1.3),
    Inches(6.2),
    Inches(1.8),
    size=16,
)

body_box(
    sl,
    "Evolution equations (Klempt et al. 2025):",
    Inches(0.3),
    Inches(2.8),
    Inches(6.2),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)
bg_rect(sl, Inches(0.3), Inches(3.3), Inches(5.9), Inches(1.5))
body_box(
    sl,
    "  φ̇ᵢ = c* φᵢ ψᵢ Σⱼ Aᵢⱼ φⱼ ψⱼ\n"
    "  ψ̇ᵢ = viability dynamics\n"
    "  Free energy: Ψ = -(1/2) c* φ̄ᵀ A φ̄",
    Inches(0.4),
    Inches(3.35),
    Inches(5.7),
    Inches(1.4),
    size=15,
    color=DARK,
)

bg_rect(sl, Inches(6.5), Inches(0.8), Inches(6.5), Inches(6.4))
body_box(
    sl,
    "Symmetric A Matrix (15 parameters):",
    Inches(6.6),
    Inches(0.9),
    Inches(6.2),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)
body_box(
    sl,
    "θ = (a₁₁, a₁₂, a₂₂)        ← So–An block\n"
    "  ⊕ (a₃₃, a₃₄, a₄₄)        ← Vei–Fn block\n"
    "  ⊕ (a₁₃, a₁₄, a₂₃, a₂₄)  ← Cross terms\n"
    "  ⊕ (a₅₅)                   ← Pg self\n"
    "  ⊕ (a₁₅, a₂₅, a₃₅, a₄₅)  ← Pg cross\n\n"
    "  Aᵢⱼ > 0 : cooperative (cross-feeding)\n"
    "  Aᵢⱼ < 0 : competitive\n"
    "  Aᵢᵢ > 0 : self-growth (density-dep.)\n\n"
    "  NOTE: µᵢ (intrinsic growth) removed —\n"
    "  multiplied by α=0 (no antibiotics),\n"
    "  completely inert in forward model",
    Inches(6.6),
    Inches(1.5),
    Inches(6.2),
    Inches(5.5),
    size=15,
    color=DARK,
)

# ── Slide 4: TMCMC ───────────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Bayesian Framework: TMCMC")

bg_rect(sl, Inches(0.3), Inches(0.8), Inches(6.2), Inches(5.8))
body_box(
    sl,
    "TMCMC (Ching & Chen 2007)",
    Inches(0.4),
    Inches(0.9),
    Inches(5.9),
    Inches(0.5),
    size=17,
    bold=True,
    color=IKM,
)
body_box(
    sl,
    "Tempered sequence:\n"
    "  pₘ(θ) ∝ p(y|θ)^βₘ · p(θ),   0=β₀<…<βₘ=1\n\n"
    "Each stage m:\n"
    "  1. Resample Nₚ particles ∝ wⱼ^(m)\n"
    "  2. Estimate weighted covariance Σ̂^(m)\n"
    "  3. Mutate via MH:  θ' ~ N(θⱼ, γ²Σ̂^(m))\n\n"
    "Adaptive γ scaling:\n"
    "  init: γ² = 2.38²/15  (Roberts et al. optimal)\n"
    "  halve if ᾱ < 5%;  ×1.5 if ᾱ > 40%\n\n"
    "By-product: log-evidence ln Ẑ = Σ ln(mean wⱼ^(m))",
    Inches(0.4),
    Inches(1.5),
    Inches(5.9),
    Inches(5.0),
    size=15,
    color=DARK,
)

bg_rect(sl, Inches(6.7), Inches(0.8), Inches(6.3), Inches(2.8))
body_box(
    sl,
    "Likelihood",
    Inches(6.8),
    Inches(0.9),
    Inches(6.0),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)
body_box(
    sl,
    "ln p(y|θ) = -Σₖ,ᵢ (yᵢₖ - ŷᵢₖ(θ))² / (2σᵢ²)\n\n"
    "σᵢ = species-specific replicate noise\n"
    "     ≈ 0.05 (Fn,Pg commensal) to 0.20 (So)",
    Inches(6.8),
    Inches(1.5),
    Inches(6.0),
    Inches(2.0),
    size=15,
    color=DARK,
)

bg_rect(sl, Inches(6.7), Inches(3.8), Inches(6.3), Inches(2.8))
body_box(
    sl, "Prior", Inches(6.8), Inches(3.9), Inches(6.0), Inches(0.5), size=16, bold=True, color=IKM
)
body_box(
    sl,
    "p(θ) = Πₖ 1/(uₖ-lₖ)  (uniform, k=0..14)\n\n"
    "Bounds narrowed after Phase 1:\n"
    "  [q₀.₀₁ - 1.5σ,  q₀.₉₉ + 1.5σ]\n"
    "(clipped to original bounds)",
    Inches(6.8),
    Inches(4.5),
    Inches(6.0),
    Inches(2.0),
    size=15,
    color=DARK,
)

# ── Slide 5: Two-Phase Strategy ───────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Two-Phase Estimation Strategy")

body_box(
    sl,
    "Phase 1a  ──→  Phase 1b  ──→  Phase 2  ──→  Posterior",
    Inches(0.5),
    Inches(0.8),
    Inches(12),
    Inches(0.5),
    size=18,
    bold=True,
    color=IKM,
)

rows = [
    ("", "Phase 1a", "Phase 1b", "Phase 2"),
    ("ψᵢ", "Fixed (exp.)", "Fixed (exp.)", "Free variable"),
    ("Prior", "Wide uniform", "Narrowed (quantile)", "Phase 1b posterior"),
    ("Nₚ", "1,000", "1,000", "10,000"),
    ("Purpose", "Explore landscape", "Identify A matrix", "Full joint posterior"),
    ("Viability", "—", "—", "Soft constraint (σ=0.1)"),
]
left, top = Inches(0.3), Inches(1.4)
col_w = [Inches(2.0), Inches(2.9), Inches(2.9), Inches(2.9), Inches(2.4)]
row_h = Inches(0.65)
for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        cleft = left + sum(col_w[:ci])
        bg_c = IKM if ri == 0 else (LIGHT_BG if ri % 2 == 0 else WHITE)
        rect = sl.shapes.add_shape(1, cleft, top + ri * row_h, col_w[ci], row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_c
        rect.line.color.rgb = RGBColor(180, 200, 220)
        tb = sl.shapes.add_textbox(
            cleft + Inches(0.05),
            top + ri * row_h + Inches(0.08),
            col_w[ci] - Inches(0.1),
            row_h - Inches(0.1),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = cell
        r.font.size = Pt(14 if ri > 0 else 13)
        r.font.bold = ri == 0 or ci == 0
        r.font.color.rgb = WHITE if ri == 0 else DARK

bullet_box(
    sl,
    [
        "Phase 1 eliminates ψ-related multimodality by conditioning on experimental viability",
        "Phase 2 viability augmentation:  ln L += -Σ (ψᵢ - ψᵢᵉˣᵖ)² / (2 × 0.1²)",
        "Result: self-consistent posterior over both A and ψᵢ",
    ],
    Inches(0.3),
    Inches(5.8),
    Inches(12.7),
    Inches(1.5),
    size=15,
)

# ── Slide 6: GPU Acceleration ─────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "GPU-Parallel Forward Model: ~200× Speedup")

bg_rect(sl, Inches(0.3), Inches(0.8), Inches(6.2), Inches(5.8))
bullet_box(
    sl,
    [
        "~10⁵–10⁶ ODE solves per run",
        "(10,000 particles × 20 MH steps × ~5 stages)",
        "Implicit Euler–Newton compiled via jax.jit",
        "All Nₚ particles via jax.vmap on GPU",
        "Single fused XLA kernel (no Python overhead)",
    ],
    Inches(0.3),
    Inches(0.8),
    Inches(6.2),
    Inches(3.2),
    size=16,
    title="JAX Implementation",
)

bullet_box(
    sl,
    [
        "RW mutation avoids variable-length leapfrog",
        "(NUTS/HMC fail on GPU: warp divergence)",
        "Adaptive γ: halve if ᾱ<5%, ×1.5 if ᾱ>40%",
        "Proposal: N(θⱼ, γ² Σ̂)  clipped to prior bounds",
    ],
    Inches(0.3),
    Inches(4.0),
    Inches(6.2),
    Inches(2.8),
    size=16,
    title="Why RW-TMCMC?",
)

bg_rect(sl, Inches(6.7), Inches(0.8), Inches(6.3), Inches(5.8))
body_box(
    sl,
    "Benchmark (RTX 4090 vs 12-core CPU)",
    Inches(6.8),
    Inches(0.9),
    Inches(6.0),
    Inches(0.5),
    size=16,
    bold=True,
    color=IKM,
)

rows2 = [
    ("Metric", "CPU", "GPU", "Speedup"),
    ("Matched (Nₚ=500, DH)", "", "", ""),
    ("  Per stage", "2,400 s", "14 s", "170×"),
    ("  Total", "21,600 s", "147 s", "200×"),
    ("Production (Nₚ=10,000)", "", "", ""),
    ("  Per condition", "~120 h†", "1.7–2.6 h", "~50×"),
    ("  All 4 conditions", "~480 h†", "8 h", "~60×"),
]
for ri, row in enumerate(rows2):
    for ci, cell in enumerate(row):
        cleft = Inches(6.8) + [Inches(0), Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2)][ci]
        col_w2 = [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2)]
        bg_c2 = IKM if ri == 0 else (LIGHT_BG if ri % 2 == 0 else WHITE)
        if row[1] == "" and ri > 0:
            bg_c2 = RGBColor(210, 225, 240)
        rect = sl.shapes.add_shape(
            1, cleft, Inches(1.5) + ri * Inches(0.62), col_w2[ci], Inches(0.62)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_c2
        rect.line.color.rgb = RGBColor(180, 200, 220)
        tb = sl.shapes.add_textbox(
            cleft + Inches(0.04),
            Inches(1.55) + ri * Inches(0.62),
            col_w2[ci] - Inches(0.08),
            Inches(0.55),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = cell
        r.font.size = Pt(13)
        r.font.bold = ri == 0
        r.font.color.rgb = WHITE if ri == 0 else (IKM if "200×" in cell else DARK)

body_box(
    sl,
    "†Extrapolated from 500-particle CPU benchmark",
    Inches(6.8),
    Inches(5.9),
    Inches(6.0),
    Inches(0.5),
    size=12,
    color=RGBColor(100, 100, 100),
)

# ── Slide 7: Fit Metrics ──────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Results: Quantitative Fit Metrics")

rows3 = [
    ("Cond.", "Ph1 RMSE", "Ph1 R²(Pg)", "Ph2 RMSE", "Ph2 R²(Pg)", "ᾱ", "ln Ẑ"),
    ("CS", "0.119", "—", "0.119", "—", "0.12", "-3.3"),
    ("CH", "0.071", "—", "0.104", "—", "0.07", "-8.5"),
    ("DS", "0.020 ✓", "0.85", "0.033 ✓", "0.85", "0.13", "-2.8 ✓"),
    ("DH", "0.067", "0.92", "0.087", "0.90 ✓", "0.07", "-11.9"),
]
col_widths = [
    Inches(1.1),
    Inches(1.5),
    Inches(1.6),
    Inches(1.5),
    Inches(1.6),
    Inches(0.9),
    Inches(1.3),
]
for ri, row in enumerate(rows3):
    for ci, cell in enumerate(row):
        cleft = Inches(0.3) + sum(col_widths[:ci])
        bg_c3 = IKM if ri == 0 else (LIGHT_BG if ri % 2 == 1 else WHITE)
        if "✓" in cell:
            bg_c3 = RGBColor(200, 230, 200)
        rect = sl.shapes.add_shape(
            1, cleft, Inches(0.85) + ri * Inches(0.65), col_widths[ci], Inches(0.65)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_c3
        rect.line.color.rgb = RGBColor(180, 200, 220)
        tb = sl.shapes.add_textbox(
            cleft + Inches(0.05),
            Inches(0.9) + ri * Inches(0.65),
            col_widths[ci] - Inches(0.1),
            Inches(0.58),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = cell
        r.font.size = Pt(15)
        r.font.bold = ri == 0
        r.font.color.rgb = WHITE if ri == 0 else DARK

bullet_box(
    sl,
    [
        "DS: lowest RMSE (0.033) — best identifiability among conditions",
        "Phase 2 RMSE increase ≤ 0.03 — Phase 1 structure is robust",
        "DH: strong Pg identification R²=0.90 despite free ψ",
        "Convergence: M=4–7 stages,  ᾱ=7–13%,  ESS=Nₚ=10,000",
        "Gelman-Rubin R̂ < 1.1 for all parameters (2 independent seeds)",
    ],
    Inches(0.3),
    Inches(4.25),
    Inches(12.7),
    Inches(3.0),
    size=16,
    title="Key Observations",
)

# ── Slide 8: Posterior Fits ───────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Posterior Predictive Fits")

fig_data = img_bytes(FIGS / "paper_fig2_phi_transposed.pdf", dpi=200)
add_figure(sl, fig_data, Inches(0.2), Inches(0.75), Inches(12.9), Inches(5.8))

body_box(
    sl,
    "Columns: CS, CH, DS, DH (increasing dysbiosis)  |  Rows: So, An, Vei, Fn, Pg  |  "
    "Solid: MAP φᵢ  |  Dashed: living fraction φ̄ᵢ=φᵢψᵢ  |  Shaded: 50%/80% CI (100 samples)",
    Inches(0.2),
    Inches(6.9),
    Inches(12.9),
    Inches(0.5),
    size=12,
    color=RGBColor(80, 80, 80),
)

# ── Slide 9: Heatmaps ─────────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Inferred Interaction Matrices A (MAP, Phase 2)")

fig_data = img_bytes(FIGS / "heatmap_A_4cond.png", dpi=180)
add_figure(sl, fig_data, Inches(0.2), Inches(0.75), Inches(8.5), Inches(5.8))

bullet_box(
    sl,
    [
        "CS/CH: So–An and So–Vei dominant",
        "Fn and Pg rows/cols ≈ 0  (recovered without locking)",
        "Emergent biological sparsity",
    ],
    Inches(8.9),
    Inches(1.1),
    Inches(4.2),
    Inches(2.2),
    size=15,
    title="Commensal",
)

bullet_box(
    sl,
    [
        "DS: moderate Fn–Pg activation",
        "DH: strong positive Vei→Pg and Fn→Pg",
        "â₅₅(DH)=3.43,  â₄₅(DH)=5.63 (Fn→Pg)",
        "â₄₄(DH)=4.45 (Fn self-growth)",
        "Cooperative cross-feeding drives surge",
    ],
    Inches(8.9),
    Inches(3.7),
    Inches(4.2),
    Inches(3.2),
    size=15,
    title="Dysbiotic",
)

# ── Slide 10: Posterior Distributions ─────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Posterior Parameter Distributions (Phase 2, Nₚ=10,000)")

fig_data = img_bytes(FIGS / "paper_posterior_violin.pdf", dpi=200)
add_figure(sl, fig_data, Inches(0.2), Inches(0.75), Inches(12.9), Inches(5.6))

body_box(
    sl,
    "CS/CH: Pg-related params (a₁₅,a₂₅,a₃₅,a₄₅,a₅₅) concentrate near zero WITHOUT a priori locking  |  "
    "Stars: MAP estimates  |  DS: narrowest posteriors (best identifiability)",
    Inches(0.2),
    Inches(6.85),
    Inches(12.9),
    Inches(0.55),
    size=12,
    color=RGBColor(80, 80, 80),
)

# ── Slide 11: Phase1 vs Phase2 ─────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Phase 1 vs Phase 2 Consistency Check")

fig_data = img_bytes(FIGS / "phase1_vs_phase2_map.pdf", dpi=200)
add_figure(sl, fig_data, Inches(0.2), Inches(0.75), Inches(10.0), Inches(5.0))

bullet_box(
    sl,
    [
        "CS:  r = 0.87  → fix-ψ valid",
        "CH:  r = 0.68  → acceptable",
    ],
    Inches(10.3),
    Inches(1.2),
    Inches(2.9),
    Inches(1.8),
    size=15,
    title="Commensal",
)

bullet_box(
    sl,
    [
        "DS:  r = 0.42  → Phase 2 essential",
        "DH:  r = 0.62  → largest shift in a₅₅, a₁₅, a₄₅",
        "Pg viability feeds back onto interaction params",
    ],
    Inches(10.3),
    Inches(3.2),
    Inches(2.9),
    Inches(2.8),
    size=15,
    title="Dysbiotic",
)

# ── Slide 12: Cross-condition Transferability ──────────────────────────────────
sl = add_slide(prs)
header(sl, "Cross-Condition Transferability")

rows4 = [
    ("Source \\ Target", "CS", "CH", "DS", "DH"),
    ("CS", "0.273†", "0.135★", "0.111★", "0.244"),
    ("CH", "0.229", "0.168†", "0.127★", "0.233"),
    ("DS", "0.273", "0.130", "0.062†", "0.252"),
    ("DH", "0.301", "0.129", "0.044★", "0.262†"),
]
col_w4 = [Inches(2.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)]
for ri, row in enumerate(rows4):
    for ci, cell in enumerate(row):
        cleft = Inches(0.3) + sum(col_w4[:ci])
        bg_c4 = IKM if ri == 0 else (LIGHT_BG if ri % 2 == 1 else WHITE)
        if "★" in cell:
            bg_c4 = RGBColor(200, 230, 200)
        elif "†" in cell:
            bg_c4 = RGBColor(230, 230, 230)
        rect = sl.shapes.add_shape(
            1, cleft, Inches(0.85) + ri * Inches(0.7), col_w4[ci], Inches(0.7)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_c4
        rect.line.color.rgb = RGBColor(180, 200, 220)
        tb = sl.shapes.add_textbox(
            cleft + Inches(0.05),
            Inches(0.9) + ri * Inches(0.7),
            col_w4[ci] - Inches(0.1),
            Inches(0.62),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = cell
        r.font.size = Pt(16)
        r.font.bold = ri == 0
        r.font.color.rgb = WHITE if ri == 0 else DARK

body_box(
    sl,
    "★: transfer RMSE < self-prediction   †: diagonal (self-prediction)",
    Inches(0.3),
    Inches(4.4),
    Inches(9.0),
    Inches(0.5),
    size=13,
    color=RGBColor(80, 80, 80),
)

bullet_box(
    sl,
    [
        "DH→DS = 0.044  BEATS DS self (0.062) → shared dysbiotic backbone",
        "CS→CH = 0.135  vs. CH self 0.168 → shared commensal backbone",
        "Cross-state: DS→CS = 0.273, CS→DH = 0.244 (large degradation)",
        "ρ(DH-CS) = -0.49 (anticorrelated) → topological reorganisation",
        "Health-to-disease ≠ parameter rescaling",
    ],
    Inches(9.5),
    Inches(0.85),
    Inches(3.6),
    Inches(4.5),
    size=15,
    title="Interpretation",
)

# ── Slide 13: Independent Validation (1/2) ────────────────────────────────────
sl = add_slide(prs)
header(sl, "Independent Validation (not used in calibration) (1/2)")

items_v = [
    (
        "1. pH prediction  (r = 0.84, R² = 0.71, N=12)",
        "Post-hoc: pH = 6.95 − 0.74 φ_So − 0.48 φ_Vei\n"
        "Negative So coefficient: lactate-driven acidification ✓",
    ),
    (
        "2. Gingipain–Pg correlation  (r = 0.90, 5 timepoints)",
        "Both show delayed onset at Day 15–21\n" "Gingipain = principal virulence factor of Pg ✓",
    ),
]
for i, (title_v, detail) in enumerate(items_v):
    bg_rect(sl, Inches(0.3), Inches(0.95) + i * Inches(2.85), Inches(12.7), Inches(2.65))
    body_box(
        sl,
        title_v,
        Inches(0.4),
        Inches(1.05) + i * Inches(2.85),
        Inches(12.4),
        Inches(0.6),
        size=18,
        bold=True,
        color=IKM,
    )
    body_box(
        sl,
        detail,
        Inches(0.5),
        Inches(1.75) + i * Inches(2.85),
        Inches(12.2),
        Inches(1.6),
        size=16,
        color=DARK,
    )

# ── Slide 14: Independent Validation (2/2) ────────────────────────────────────
sl = add_slide(prs)
header(sl, "Independent Validation (not used in calibration) (2/2)")

items_v = [
    (
        "3. Metabolic sign consistency",
        "A₁₂ = +0.76 (An→Vei): lactate/succinate cross-feeding ✓\n"
        "A₁₃ = −2.73 (So→Vei): niche competition dominates ✓",
    ),
    (
        "4. Growth rate parity",
        "Planktonic doubling: 1.68 h (commensal) ≈ 1.64 h (dysbiotic)\n"
        "Condition differences arise from interactions, not fitness ✓",
    ),
]
for i, (title_v, detail) in enumerate(items_v):
    bg_rect(sl, Inches(0.3), Inches(0.95) + i * Inches(2.85), Inches(12.7), Inches(2.65))
    body_box(
        sl,
        title_v,
        Inches(0.4),
        Inches(1.05) + i * Inches(2.85),
        Inches(12.4),
        Inches(0.6),
        size=18,
        bold=True,
        color=IKM,
    )
    body_box(
        sl,
        detail,
        Inches(0.5),
        Inches(1.75) + i * Inches(2.85),
        Inches(12.2),
        Inches(1.6),
        size=16,
        color=DARK,
    )

# ── Slide 15: Discussion ──────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Discussion: Biological Interpretation")

bullet_box(
    sl,
    [
        "So–An co-aggregation + So–Vei lactate exchange: stable mutualistic core",
        "HOBIC flow shifts So abundance; community structure preserved",
        "Fn/Pg entries near zero (posterior recovery, not locking)",
    ],
    Inches(0.3),
    Inches(0.8),
    Inches(6.2),
    Inches(2.5),
    size=16,
    title="Commensal States (CS, CH)",
)

bullet_box(
    sl,
    [
        "Vei→Pg (pH) + Fn→Pg (peptides) strongly activated",
        "ρ(DH–CS) = −0.49: topological reorganisation",
        "Cannot be explained by parameter rescaling",
    ],
    Inches(0.3),
    Inches(3.5),
    Inches(6.2),
    Inches(2.8),
    size=16,
    title="Dysbiotic Transition",
)

bg_rect(sl, Inches(6.7), Inches(0.8), Inches(6.3), Inches(5.8))
body_box(
    sl,
    "Model Parsimony: Effective Dimensionality",
    Inches(6.8),
    Inches(0.9),
    Inches(6.0),
    Inches(0.5),
    size=18,
    bold=True,
    color=IKM,
)
bullet_box(
    sl,
    [
        "Metric: rᵢ = σ_post,ᵢ / Δ_prior  (Δ_prior=6)",
        "All 15 params: rᵢ < 0.43 in every condition",
        "DS: tightest posteriors (mean r = 0.07)",
        "Pg params in commensal: rᵢ ≤ 0.10",
        "Interpretation: data-informed sparsity → full 15-param model justified",
    ],
    Inches(6.8),
    Inches(1.5),
    Inches(6.0),
    Inches(3.3),
    size=16,
)
body_box(
    sl,
    "Acceptance rates 7–13% (< optimal 23% for d=15) reflect non-Gaussianity;\n"
    "high ESS confirms mixing",
    Inches(6.8),
    Inches(4.95),
    Inches(6.0),
    Inches(1.0),
    size=13,
    color=MID_GRAY,
)

# ── Slide 16: Conclusions ─────────────────────────────────────────────────────
sl = add_slide(prs)
header(sl, "Conclusions")

conclusions = [
    (
        "Two-Phase Full-Discovery Estimation",
        "Phase 1 (fix-ψ): rapid A-matrix identification via narrowed priors\n"
        "Phase 2 (free ψ): physically consistent joint posterior\n"
        "All 15 params estimated without a priori locking → emergent sparsity",
    ),
    (
        "GPU-Accelerated TMCMC: ~200× Speedup",
        "JAX + jax.vmap on RTX 4090; Nₚ=10,000 runs in ~2 h/condition\n"
        "RW-TMCMC avoids warp divergence (NUTS/HMC not GPU-friendly)",
    ),
    (
        "Validated Across 4 Conditions + Independent Data",
        "pH prediction r=0.84, gingipain r=0.90, metabolic sign consistent\n"
        "Within-state MAP transfer outperforms cross-state",
    ),
    (
        "Health-to-Disease = Topological Reorganisation",
        "ρ(DH–CS)=−0.49; within-state backbone generalises across cultivation modes\n"
        "Not recoverable by parameter rescaling",
    ),
    (
        "Future Work",
        "Surrogate acceleration (DeepONet ~100× additional speedup)\n"
        "Spatial PDE extension, VEM-based FEM, Siddiqui 2021 (6-species) validation",
    ),
]
for i, (ttl, det) in enumerate(conclusions):
    bg_col = LIGHT_BG if i < 4 else RGBColor(230, 240, 220)
    bg_rect(
        sl, Inches(0.3), Inches(0.8) + i * Inches(1.28), Inches(12.7), Inches(1.22), color=bg_col
    )
    num_color = IKM if i < 4 else RGBColor(60, 130, 60)
    body_box(
        sl,
        f"{i+1}. {ttl}",
        Inches(0.4),
        Inches(0.85) + i * Inches(1.28),
        Inches(12.4),
        Inches(0.45),
        size=16,
        bold=True,
        color=num_color,
    )
    body_box(
        sl,
        det,
        Inches(0.5),
        Inches(1.3) + i * Inches(1.28),
        Inches(12.2),
        Inches(0.7),
        size=14,
        color=DARK,
    )

# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
