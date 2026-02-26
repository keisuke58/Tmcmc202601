#!/usr/bin/env python3
"""
Generate polished PowerPoint presentation for
Biologically-Constrained Parameter Reduction for 5-Species Biofilm Model.
Full content: methods, algorithms, equations, all 4-condition results, diagnostics.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ────────────────────────────────────────────────────────────────────
BASE = "/home/nishioka/IKM_Hiwi/Tmcmc202601/data_5species"
RUNS = os.path.join(BASE, "_runs")
EXPFIG = os.path.join(BASE, "experiment_fig")
DOCS = os.path.join(BASE, "docs")
EQIMG = os.path.join(DOCS, "_eq_img")

CONDS = {
    "Commensal_Static": {
        "dir": os.path.join(RUNS, "Commensal_Static_20260208_002100"),
        "label": "Commensal Static",
        "label_ja": "健康・静的",
        "short": "CS",
        "desc_en": (
            "Healthy state under static culture. Nutrient-limited environment "
            "leads to stable, low biomass. Strict parameter locks (N_locked = 9) "
            "enforce pathogen suppression, reflecting the clinically observed "
            "stable symbiotic community."
        ),
        "desc_ja": (
            "静的培養下の健康な状態。栄養制限環境により安定的な低バイオマスとなる。"
            "厳格なパラメータロック (N_locked = 9) により病原菌の増殖が抑制され、"
            "臨床的に観察される安定した共生群集を反映。"
        ),
        "locked": 9,
        "finding_en": (
            "Strict locks prevent pathogen growth. The model accurately reproduces "
            "the stable commensal state observed in static experiments. Early "
            "colonizers (S.o, A.n) maintain dominance throughout."
        ),
        "finding_ja": (
            "厳格なロックにより病原菌の増殖が防がれ、静的実験で観察された"
            "安定した共生状態を正確に再現。初期定着菌(S.o, A.n)が終始優位を維持。"
        ),
        "color": RGBColor(0x28, 0xA7, 0x45),  # green
    },
    "Dysbiotic_Static": {
        "dir": os.path.join(RUNS, "Dysbiotic_Static_20260207_203752"),
        "label": "Dysbiotic Static",
        "label_ja": "疾患・静的",
        "short": "DS",
        "desc_en": (
            "Disease state under static culture. Pathogens are present but limited "
            "by metabolite accumulation. Standard biological locks (N_locked = 5) "
            "are applied."
        ),
        "desc_ja": (
            "静的培養下の疾患状態。病原菌は存在するが代謝産物の蓄積により制限。"
            "標準的な生物学的ロック (N_locked = 5) を適用。"
        ),
        "locked": 5,
        "finding_en": (
            "Pathogen interactions are estimated but their magnitude is smaller "
            "compared to HOBIC conditions. This confirms that saliva flow is "
            "essential for full dysbiosis progression."
        ),
        "finding_ja": (
            "病原菌の相互作用は推定されたがHOBIC条件と比較して規模が小さく、"
            "完全なディスバイオシスの進行には流れ(Flow)が不可欠であることを裏付け。"
        ),
        "color": RGBColor(0xFF, 0x8C, 0x00),  # orange
    },
    "Commensal_HOBIC": {
        "dir": os.path.join(RUNS, "Commensal_HOBIC_20260208_002100"),
        "label": "Commensal HOBIC",
        "label_ja": "健康・流動 (HOBIC)",
        "short": "CH",
        "desc_en": (
            "Healthy state with saliva flow (HOBIC). Characterized by S. oralis "
            "'Blue Bloom' — high proliferation of early colonizers under flow. "
            "Strict locks (N_locked = 8) allow S. oralis growth while suppressing "
            "pathogens."
        ),
        "desc_ja": (
            "流れのある健康な状態 (HOBIC)。S. oralisの高い増殖 ('Blue Bloom') が特徴。"
            "厳格なロック (N_locked = 8) がS. oralisの増殖を許容しつつ病原菌を抑制。"
        ),
        "locked": 8,
        "finding_en": (
            "The model correctly identifies early colonizer dominance while "
            "suppressing pathogen populations. The 'Blue Bloom' observation "
            "is quantitatively reproduced by the MAP estimate."
        ),
        "finding_ja": (
            "初期定着菌の優位性を正しく特定しつつ病原菌集団を低く抑制。"
            "'Blue Bloom' の観察結果をMAP推定値により定量的に再現。"
        ),
        "color": RGBColor(0x00, 0x72, 0xCE),  # blue
    },
    "Dysbiotic_HOBIC": {
        "dir": os.path.join(RUNS, "Dysbiotic_HOBIC_20260208_002100"),
        "label": "Dysbiotic HOBIC (Surge)",
        "label_ja": "疾患・流動 (Surge)",
        "short": "DH",
        "desc_en": (
            "Disease state with saliva flow (HOBIC). Characterized by the "
            "explosive 'Surge' of F. nucleatum and P. gingivalis. All parameter "
            "locks are released (N_locked = 0, Discovery Mode) to capture "
            "complex cross-feeding dynamics."
        ),
        "desc_ja": (
            "流れのある疾患状態 (HOBIC)。F. nucleatumとP. gingivalisの爆発的な "
            "'Surge (急増)' が特徴。複雑なクロスフィーディング動態を捕捉するため"
            "全パラメータロックを解除 (N_locked = 0, Discovery Mode)。"
        ),
        "locked": 0,
        "finding_en": (
            "Unlock All (Discovery Mode) successfully captures the non-linear "
            "Surge via strong Vei–P.g positive feedback. TMCMC identifies "
            "the critical cross-feeding loop (Index 18) as the driver of "
            "explosive P.g growth."
        ),
        "finding_ja": (
            "全ロック解除 (Discovery Mode) により、Vei–P.g間の強い正のフィードバック"
            "を介した非線形な急増 (Surge) を再現に成功。TMCMCがP.gの爆発的増殖の"
            "駆動因子であるクロスフィーディングループ (Index 18) を特定。"
        ),
        "color": RGBColor(0xE8, 0x3E, 0x3E),  # red
    },
}

# ── Color palette ────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DARK_NAVY = RGBColor(0x08, 0x10, 0x1C)
TITLE_BLUE = RGBColor(0x00, 0x72, 0xCE)
ACCENT_RED = RGBColor(0xE8, 0x3E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG = RGBColor(0xF0, 0xF2, 0xF5)
SECTION_BG = RGBColor(0x00, 0x3D, 0x7A)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
BORDER_BLUE = RGBColor(0x00, 0x56, 0xA3)
SOFT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
ALGO_BG = RGBColor(0xFD, 0xFD, 0xFD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
slide_counter = [0]


# ── Helper functions ─────────────────────────────────────────────────────────
def _bg(slide, color=SLIDE_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color, border_color=None, border_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_width or 1)
    else:
        s.line.fill.background()
    return s


def _rounded_rect(slide, l, t, w, h, color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def _txt(
    slide,
    l,
    t,
    w,
    h,
    text,
    sz=18,
    bold=False,
    color=DARK_TEXT,
    align=PP_ALIGN.LEFT,
    font="Calibri",
    valign=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = valign
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox


def _multi(
    slide,
    l,
    t,
    w,
    h,
    lines,
    default_sz=16,
    default_color=DARK_TEXT,
    default_bold=False,
    font="Calibri",
    align=PP_ALIGN.LEFT,
):
    """lines: list of str or tuple (text, sz, bold, color)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, b, c = line, default_sz, default_bold, default_color
        else:
            txt = line[0]
            sz = line[1] if len(line) > 1 else default_sz
            b = line[2] if len(line) > 2 else default_bold
            c = line[3] if len(line) > 3 else default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def _img(slide, path, l, t, width=None, height=None):
    if os.path.exists(path):
        kw = {"image_file": path, "left": l, "top": t}
        if width:
            kw["width"] = width
        if height:
            kw["height"] = height
        return slide.shapes.add_picture(**kw)
    else:
        _txt(
            slide,
            l,
            t,
            Inches(4),
            Inches(0.4),
            f"[Missing: {os.path.basename(path)}]",
            sz=11,
            color=ACCENT_RED,
        )
        return None


def _footer(slide, text="Nishioka  |  5-Species Biofilm  |  TMCMC  |  IKM, LUH"):
    slide_counter[0] += 1
    _rect(slide, Inches(0), SH - Inches(0.38), SW, Inches(0.38), NAVY)
    _txt(
        slide, Inches(0.4), SH - Inches(0.36), Inches(8), Inches(0.32), text, sz=9, color=LIGHT_GRAY
    )
    _txt(
        slide,
        SW - Inches(1.2),
        SH - Inches(0.36),
        Inches(1),
        Inches(0.32),
        str(slide_counter[0]),
        sz=9,
        color=LIGHT_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def _header(slide, title, subtitle=""):
    # Top accent bar
    _rect(slide, Inches(0), Inches(0), SW, Inches(0.06), TITLE_BLUE)
    _txt(
        slide,
        Inches(0.6),
        Inches(0.2),
        Inches(12),
        Inches(0.7),
        title,
        sz=30,
        bold=True,
        color=TITLE_BLUE,
    )
    if subtitle:
        _txt(
            slide,
            Inches(0.6),
            Inches(0.78),
            Inches(12),
            Inches(0.35),
            subtitle,
            sz=16,
            color=MID_GRAY,
        )


def _section(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, SECTION_BG)
    # Decorative line
    _rect(s, Inches(1.5), Inches(2.2), Inches(10), Inches(0.04), GOLD)
    _txt(
        s,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.5),
        title,
        sz=44,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    if subtitle:
        _txt(
            s,
            Inches(1),
            Inches(4.3),
            Inches(11),
            Inches(0.8),
            subtitle,
            sz=22,
            color=LIGHT_GRAY,
            align=PP_ALIGN.CENTER,
        )
    _rect(s, Inches(1.5), Inches(5.2), Inches(10), Inches(0.04), GOLD)
    _footer(s)
    return s


def _card(slide, l, t, w, h, border_color=None):
    return _rounded_rect(
        slide, l, t, w, h, CARD_BG, border_color=border_color or RGBColor(0xDD, 0xDD, 0xDD)
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s, NAVY)
_rect(s, Inches(0), Inches(0), SW, Inches(0.08), TITLE_BLUE)
_rect(s, Inches(0), SH - Inches(0.08), SW, Inches(0.08), TITLE_BLUE)
_txt(
    s,
    Inches(1),
    Inches(1.2),
    Inches(11.3),
    Inches(1.5),
    "Biologically-Constrained\nParameter Reduction",
    sz=46,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(1),
    Inches(3.0),
    Inches(11.3),
    Inches(0.7),
    "for 5-Species Oral Biofilm Model",
    sz=32,
    bold=False,
    color=GOLD,
    align=PP_ALIGN.CENTER,
)
_rect(s, Inches(4), Inches(3.9), Inches(5.3), Inches(0.03), TITLE_BLUE)
_txt(
    s,
    Inches(1),
    Inches(4.2),
    Inches(11.3),
    Inches(0.6),
    "生物学的制約に基づく5菌種バイオフィルムモデルのパラメータ削減",
    sz=18,
    color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(1),
    Inches(5.0),
    Inches(11.3),
    Inches(0.6),
    "Comprehensive Bayesian Uncertainty Quantification\nunder Commensal & Dysbiotic Conditions",
    sz=18,
    color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(1),
    Inches(6.2),
    Inches(11.3),
    Inches(0.5),
    "Keisuke Nishioka  \u2502  Institut f\u00fcr Kontinuumsmechanik, Leibniz Universit\u00e4t Hannover  \u2502  February 2026",
    sz=14,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: OUTLINE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Outline / 目次")
items_l = [
    ("\u2460  Introduction & Biological Background", 19, True, DARK_TEXT),
    ("      はじめに・生物学的背景", 14, False, MID_GRAY),
    ("\u2461  5-Species Biofilm Model", 19, True, DARK_TEXT),
    ("      5菌種バイオフィルムモデル", 14, False, MID_GRAY),
    ("\u2462  Biologically-Constrained Parameter Reduction", 19, True, DARK_TEXT),
    ("      生物学的制約パラメータ削減", 14, False, MID_GRAY),
    ("\u2463  Bayesian Inference Framework", 19, True, DARK_TEXT),
    ("      ベイズ推定の枠組み", 14, False, MID_GRAY),
    ("\u2464  TMCMC Algorithm", 19, True, DARK_TEXT),
    ("      TMCMCアルゴリズム", 14, False, MID_GRAY),
    ("\u2465  4-Stage Sequential Estimation", 19, True, DARK_TEXT),
    ("      4段階逐次推定", 14, False, MID_GRAY),
]
items_r = [
    ("\u2466  Results: Commensal Static", 19, True, GREEN),
    ("      結果：健康・静的", 14, False, MID_GRAY),
    ("\u2467  Results: Dysbiotic Static", 19, True, ORANGE),
    ("      結果：疾患・静的", 14, False, MID_GRAY),
    ("\u2468  Results: Commensal HOBIC", 19, True, TITLE_BLUE),
    ("      結果：健康・流動", 14, False, MID_GRAY),
    ("\u2469  Results: Dysbiotic HOBIC (Surge)", 19, True, ACCENT_RED),
    ("      結果：疾患・流動 (Surge)", 14, False, MID_GRAY),
    ("\u246A  Comparative Analysis", 19, True, DARK_TEXT),
    ("      比較分析", 14, False, MID_GRAY),
    ("\u246B  Conclusion & Future Work", 19, True, DARK_TEXT),
    ("      結論・今後の課題", 14, False, MID_GRAY),
]
_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.6))
_multi(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3), items_l)
_card(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.6))
_multi(s, Inches(7.1), Inches(1.5), Inches(5.8), Inches(5.3), items_r)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1: INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2460 Introduction", "はじめに・生物学的背景")

# -- Background slide --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(
    s,
    "Biological Background",
    "Peri-implantitis & Multi-species Biofilm / インプラント周囲炎と多菌種バイオフィルム",
)

_card(s, Inches(0.4), Inches(1.3), Inches(7.3), Inches(5.5))
lines = [
    ("Peri-implantitis", 22, True, ACCENT_RED),
    ("Inflammatory disease around dental implants caused by", 16),
    ("dysbiotic shift in the oral microbiome.", 16),
    ("歯科インプラント周囲の炎症性疾患。口腔マイクロバイオームの", 14, False, MID_GRAY),
    ("ディスバイオティックシフトにより引き起こされる。", 14, False, MID_GRAY),
    ("", 8),
    ("5 Key Bacterial Species (Heine et al.)", 20, True, TITLE_BLUE),
    ("", 4),
    ("Sp 0: Streptococcus oralis (S.o)", 16, True, DARK_TEXT),
    ("        Early colonizer — initiates biofilm formation", 14, False, MID_GRAY),
    ("Sp 1: Actinomyces naeslundii (A.n)", 16, True, DARK_TEXT),
    ("        Early colonizer — co-aggregation partner", 14, False, MID_GRAY),
    ("Sp 2: Veillonella spp. (Vei)", 16, True, DARK_TEXT),
    ("        Metabolic bridge — lactate consumer", 14, False, MID_GRAY),
    ("Sp 3: Fusobacterium nucleatum (F.n)", 16, True, DARK_TEXT),
    ("        Bridging organism — connects early & late colonizers", 14, False, MID_GRAY),
    ("Sp 4: Porphyromonas gingivalis (P.g)", 16, True, ACCENT_RED),
    ("        Late colonizer (pathogen) — key in dysbiosis", 14, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(6.8), Inches(5.2), lines)

# Reference paper figure
fig_ref = os.path.join(EXPFIG, "froh-06-1649419-g001.jpg")
_img(s, fig_ref, Inches(8.0), Inches(1.3), width=Inches(4.8))
_txt(
    s,
    Inches(8.0),
    Inches(6.3),
    Inches(4.8),
    Inches(0.4),
    "Heine et al. (2021) — Biofilm architecture",
    sz=11,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_footer(s)

# -- Interaction network slide --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Interaction Network (Figure 4C)", "相互作用ネットワーク")

_card(s, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.5))
lines = [
    ("Known Species Interactions", 20, True, TITLE_BLUE),
    ("", 6),
    ("Based on co-culture experiments by Heine et al.,", 16),
    ("the following interactions are established:", 16),
    ("", 8),
    ("\u2714  S.o \u2194 A.n   (co-aggregation)", 16, False, GREEN),
    ("\u2714  S.o \u2194 Vei   (lactate cross-feeding)", 16, False, GREEN),
    ("\u2714  S.o \u2194 F.n   (co-aggregation)", 16, False, GREEN),
    ("\u2714  A.n \u2194 P.g   — NO interaction", 16, False, ACCENT_RED),
    ("\u2714  Vei \u2194 F.n   — NO interaction", 16, False, ACCENT_RED),
    ("\u2714  A.n \u2194 Vei   — NO interaction", 16, False, ACCENT_RED),
    ("\u2714  A.n \u2194 F.n   — NO interaction", 16, False, ACCENT_RED),
    ("\u2714  S.o \u2194 P.g   — NO interaction", 16, False, ACCENT_RED),
    ("", 8),
    ("\u2192 5 pairs locked to zero", 18, True, ACCENT_RED),
    ("\u2192 Free parameters: 20 \u2192 15", 18, True, GREEN),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(5.2), Inches(5.2), lines)

fig4c = os.path.join(EXPFIG, "fig4C_reproduced_v2.png")
_img(s, fig4c, Inches(6.2), Inches(1.2), width=Inches(6.5), height=Inches(5.2))
_txt(
    s,
    Inches(6.2),
    Inches(6.5),
    Inches(6.5),
    Inches(0.3),
    "Reproduced interaction network from Heine et al. Figure 4C",
    sz=11,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2: 5-SPECIES MODEL
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2461 5-Species Biofilm Model", "5菌種バイオフィルムモデル")

s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Governing Equations", "支配方程式 — Extended Hamilton Principle")

# Left card: volume fraction
_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(2.8), TITLE_BLUE)
_txt(
    s,
    Inches(0.7),
    Inches(1.4),
    Inches(5.8),
    Inches(0.4),
    "Volume Fraction Dynamics  \u03c6\u1d62(t)",
    sz=20,
    bold=True,
    color=TITLE_BLUE,
)
_txt(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(0.3), "体積分率の動態", sz=14, color=MID_GRAY)
_img(s, os.path.join(EQIMG, "eq_phi.png"), Inches(0.5), Inches(2.15), width=Inches(5.9))
lines = [
    ("\u03c6\u1d62 : volume fraction of species i", 14, False, MID_GRAY),
    ("r\u1d62, d\u1d62 : growth and death rates", 14, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(3.4), Inches(5.8), Inches(0.6), lines)

# Right card: survival fraction
_card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(2.8), TITLE_BLUE)
_txt(
    s,
    Inches(7.2),
    Inches(1.4),
    Inches(5.6),
    Inches(0.4),
    "Survival Fraction Dynamics  \u03c8\u1d62(t)",
    sz=20,
    bold=True,
    color=TITLE_BLUE,
)
_txt(s, Inches(7.2), Inches(1.8), Inches(5.6), Inches(0.3), "生存分率の動態", sz=14, color=MID_GRAY)
_img(s, os.path.join(EQIMG, "eq_psi.png"), Inches(7.0), Inches(2.15), width=Inches(5.9))
lines = [
    ("\u03c8\u1d62 : survival fraction of species i", 14, False, MID_GRAY),
    (
        "a\u1d62\u2c7c : interaction coefficient (j \u2192 i), K : half-saturation",
        14,
        False,
        MID_GRAY,
    ),
]
_multi(s, Inches(7.2), Inches(3.4), Inches(5.6), Inches(0.6), lines)

# Bottom card: parameter summary
_card(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.5), DARK_TEXT)
lines = [
    ("Parameter Space Overview / パラメータ空間の概要", 20, True, DARK_TEXT),
    ("", 6),
    ("5\u00d75 Interaction matrix A : 25 entries", 16),
    (
        "  \u2192 Symmetry A\u1d62\u2c7c = A\u2c7c\u1d62 : 15 unique entries (5 diagonal + 10 off-diagonal)",
        16,
        False,
        TITLE_BLUE,
    ),
    ("  \u2192 Lock 5 absent interactions : 10 free interaction parameters", 16, False, GREEN),
    ("5 Decay rates b = [b\u2081, b\u2082, b\u2083, b\u2084, b\u2085] : always free", 16),
    ("", 6),
    ("", 4),
]
_multi(s, Inches(0.7), Inches(4.6), Inches(12), Inches(1.5), lines)
_img(s, os.path.join(EQIMG, "eq_reduction.png"), Inches(1.5), Inches(6.0), width=Inches(10))
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3: PARAMETER REDUCTION
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2462 Parameter Reduction", "生物学的制約に基づくパラメータ削減")

# -- Symmetric matrix & mapping --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Symmetric Interaction Matrix", "対称相互作用行列とパラメータ写像")

_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.5))
_txt(
    s,
    Inches(0.7),
    Inches(1.4),
    Inches(5.8),
    Inches(0.4),
    "5\u00d75 Symmetric Matrix A  (A\u1d62\u2c7c = A\u2c7c\u1d62)",
    sz=18,
    bold=True,
    color=TITLE_BLUE,
)
_img(s, os.path.join(EQIMG, "eq_matrix_A.png"), Inches(0.5), Inches(1.8), width=Inches(5.8))
_img(s, os.path.join(EQIMG, "eq_decay.png"), Inches(0.8), Inches(4.6), width=Inches(5.0))
lines = [
    ("Locked Indices (set to 0):", 15, True, ACCENT_RED),
    ("  Idx 6: Vei\u2194F.n   Idx 12: A.n\u2194Vei   Idx 13: A.n\u2194F.n", 13, False, ACCENT_RED),
    ("  Idx 16: S.o\u2194P.g   Idx 17: A.n\u2194P.g", 13, False, ACCENT_RED),
]
_multi(s, Inches(0.7), Inches(5.4), Inches(5.8), Inches(1.2), lines)

# Right: condition-specific lock rules
_card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(5.5))
lines = [
    ("Condition-Specific Lock Rules", 18, True, TITLE_BLUE),
    ("条件別ロックルール", 14, False, MID_GRAY),
    ("", 8),
    ("1. Commensal Static  (N_locked = 9)", 16, True, GREEN),
    ("   Standard 5 + pathogen interactions locked", 14),
    (
        "   Additional: a\u2083\u2085(F.n\u2194P.g), a\u2084\u2085(Vei\u2194P.g),",
        13,
        False,
        MID_GRAY,
    ),
    ("               b\u2084(P.g decay), b\u2083(F.n decay)", 13, False, MID_GRAY),
    ("", 6),
    ("2. Dysbiotic Static  (N_locked = 5)", 16, True, ORANGE),
    ("   Standard biological locks only", 14),
    ("   Base 5 absent interactions", 13, False, MID_GRAY),
    ("", 6),
    ("3. Commensal HOBIC  (N_locked = 8)", 16, True, TITLE_BLUE),
    ("   Strict locks, S. oralis growth allowed", 14),
    ("   Similar to CS but HOBIC-adapted", 13, False, MID_GRAY),
    ("", 6),
    ("4. Dysbiotic HOBIC  (N_locked = 0)", 16, True, ACCENT_RED),
    ("   \u2605 UNLOCK ALL \u2014 Discovery Mode", 14, True, ACCENT_RED),
    ("   All 20 parameters free for Surge detection", 13, False, MID_GRAY),
]
_multi(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(5.2), lines)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 4: BAYESIAN INFERENCE
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2463 Bayesian Inference Framework", "ベイズ推定の枠組み")

s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Bayesian Inverse Problem", "ベイズ逆問題")

# Bayes theorem card
_card(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(2.2), TITLE_BLUE)
_txt(
    s,
    Inches(0.8),
    Inches(1.4),
    Inches(12),
    Inches(0.4),
    "Bayes' Theorem / ベイズの定理",
    sz=22,
    bold=True,
    color=TITLE_BLUE,
    align=PP_ALIGN.CENTER,
)
_img(s, os.path.join(EQIMG, "eq_bayes_labels.png"), Inches(2.5), Inches(1.8), width=Inches(8))
_img(s, os.path.join(EQIMG, "eq_bayes_components.png"), Inches(3), Inches(2.8), width=Inches(7))

# Three sub-cards with equation images
# Likelihood
_card(s, Inches(0.4), Inches(3.8), Inches(4), Inches(3.0), TITLE_BLUE)
_txt(
    s,
    Inches(0.7),
    Inches(3.9),
    Inches(3.5),
    Inches(0.35),
    "Likelihood / 尤度関数",
    sz=18,
    bold=True,
    color=TITLE_BLUE,
)
_txt(
    s,
    Inches(0.7),
    Inches(4.3),
    Inches(3.5),
    Inches(0.25),
    "Gaussian noise model:",
    sz=13,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_likelihood.png"), Inches(0.5), Inches(4.6), width=Inches(3.8))
lines = [
    ("d\u1d62\u2c7c : observed data", 12, False, MID_GRAY),
    ("g\u1d62(t\u2c7c,\u03b8) : model prediction", 12, False, MID_GRAY),
    ("\u03c3\u00b2 : noise variance", 12, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(5.6), Inches(3.5), Inches(1.0), lines)

# Prior
_card(s, Inches(4.7), Inches(3.8), Inches(4), Inches(3.0), GREEN)
_txt(
    s,
    Inches(5.0),
    Inches(3.9),
    Inches(3.5),
    Inches(0.35),
    "Prior / 事前分布",
    sz=18,
    bold=True,
    color=GREEN,
)
_txt(
    s,
    Inches(5.0),
    Inches(4.3),
    Inches(3.5),
    Inches(0.25),
    "With biological constraints:",
    sz=13,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_prior.png"), Inches(4.8), Inches(4.6), width=Inches(3.8))
_img(s, os.path.join(EQIMG, "eq_prior_uniform.png"), Inches(5.0), Inches(5.5), width=Inches(3.2))
_txt(
    s,
    Inches(5.0),
    Inches(6.2),
    Inches(3.5),
    Inches(0.4),
    "\u03b4(\u03b8_k) : Dirac delta for locked params\nL = set of locked indices",
    sz=11,
    color=MID_GRAY,
)

# Evidence
_card(s, Inches(9.0), Inches(3.8), Inches(3.9), Inches(3.0), ORANGE)
_txt(
    s,
    Inches(9.3),
    Inches(3.9),
    Inches(3.5),
    Inches(0.35),
    "Evidence / エビデンス",
    sz=18,
    bold=True,
    color=ORANGE,
)
_txt(
    s,
    Inches(9.3),
    Inches(4.3),
    Inches(3.5),
    Inches(0.25),
    "Marginal likelihood:",
    sz=13,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_marginal.png"), Inches(9.1), Inches(4.6), width=Inches(3.7))
lines = [
    ("Estimated as by-product", 12, False, MID_GRAY),
    ("of TMCMC algorithm.", 12, False, MID_GRAY),
    ("Used for Bayes factor", 12, False, MID_GRAY),
    ("model comparison.", 12, False, MID_GRAY),
]
_multi(s, Inches(9.3), Inches(5.6), Inches(3.3), Inches(1.0), lines)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 5: TMCMC ALGORITHM
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2464 TMCMC Algorithm", "Transitional Markov Chain Monte Carlo")

# -- Overview slide --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "TMCMC: Algorithm Overview", "TMCMCアルゴリズムの概要")

_card(s, Inches(0.4), Inches(1.3), Inches(7.5), Inches(5.5))
lines = [
    ("Core Idea: Adaptive Tempering", 22, True, TITLE_BLUE),
    ("核心: 適応的テンパリング", 14, False, MID_GRAY),
    ("", 4),
    ("Tempered distribution sequence:", 18, True, DARK_TEXT),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(7.0), Inches(1.2), lines)
_img(s, os.path.join(EQIMG, "eq_tempered.png"), Inches(0.8), Inches(2.55), width=Inches(6.5))
_img(s, os.path.join(EQIMG, "eq_beta_seq.png"), Inches(1.2), Inches(3.35), width=Inches(5.5))
lines = [
    ("\u03b2 = 0 : samples from prior p(\u03b8)", 15, False, GREEN),
    ("\u03b2 = 1 : samples from full posterior p(\u03b8|D)", 15, False, ACCENT_RED),
    ("", 4),
    ("Adaptive \u03b2 selection:  CoV target:", 16, True, DARK_TEXT),
]
_multi(s, Inches(0.7), Inches(4.0), Inches(7.0), Inches(1.0), lines)
_img(s, os.path.join(EQIMG, "eq_cov_target.png"), Inches(0.8), Inches(5.0), width=Inches(6.8))
lines = [
    ("Advantages over standard MCMC:", 16, True, DARK_TEXT),
    ("  \u2022 No burn-in  \u2022 Multimodal  \u2022 Evidence  \u2022 Parallel", 14),
]
_multi(s, Inches(0.7), Inches(5.7), Inches(7.0), Inches(0.9), lines)

# Diagram: flow from prior to posterior
_card(s, Inches(8.2), Inches(1.3), Inches(4.8), Inches(5.5))
lines = [
    ("TMCMC Stage Flow", 18, True, TITLE_BLUE),
    ("", 8),
    (
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
        14,
        False,
        TITLE_BLUE,
    ),
    ("\u2502  Stage 0: \u03b2=0     \u2502", 14, False, TITLE_BLUE),
    ("\u2502  Prior samples    \u2502", 14, False, TITLE_BLUE),
    (
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
        14,
        False,
        TITLE_BLUE,
    ),
    ("        \u2193 weight + resample + MCMC", 13, False, MID_GRAY),
    (
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
        14,
        False,
        ORANGE,
    ),
    ("\u2502  Stage 1: \u03b2=\u03b2\u2081   \u2502", 14, False, ORANGE),
    ("\u2502  Intermediate     \u2502", 14, False, ORANGE),
    (
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
        14,
        False,
        ORANGE,
    ),
    ("        \u2193 weight + resample + MCMC", 13, False, MID_GRAY),
    ("        \u22ee", 16, False, MID_GRAY),
    ("        \u2193 weight + resample + MCMC", 13, False, MID_GRAY),
    (
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
        14,
        False,
        GREEN,
    ),
    ("\u2502  Stage m: \u03b2=1.0  \u2502", 14, False, GREEN),
    ("\u2502  Full posterior   \u2502", 14, False, GREEN),
    (
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
        14,
        False,
        GREEN,
    ),
]
_multi(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.2), lines, font="Consolas")
_footer(s)


# -- Pseudocode slide --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Algorithm 1: TMCMC Procedure", "TMCMCアルゴリズムの擬似コード")

_card(s, Inches(0.4), Inches(1.3), Inches(8.5), Inches(5.7), TITLE_BLUE)
algo = [
    ("Algorithm 1: Transitional Markov Chain Monte Carlo", 16, True, TITLE_BLUE),
    ("\u2500" * 55, 10, False, LIGHT_GRAY),
    ("Input:  N particles, prior p(\u03b8), likelihood p(D|\u03b8), data D", 13, False, DARK_TEXT),
    (
        "Output: Posterior samples {\u03b8\u207d\u00b9\u207e,...,\u03b8\u207d\u1d3a\u207e}, evidence p\u0302(D)",
        13,
        False,
        DARK_TEXT,
    ),
    ("", 4),
    (" 1:  Initialize: Draw \u03b8\u207d\u2071\u207e ~ p(\u03b8) for i = 1,...,N", 13),
    (" 2:  Set \u03b2\u2080 = 0,  j = 0,  p\u0302(D) = 1", 13),
    (" 3:  while \u03b2\u2c7c < 1 do", 13),
    (" 4:      Solve for \u0394\u03b2 such that CoV(w) = 1.0", 13),
    (" 5:      \u03b2\u2c7c\u208a\u2081 = min(\u03b2\u2c7c + \u0394\u03b2, 1)", 13),
    (
        " 6:      Compute weights: w\u207d\u2071\u207e = p(D|\u03b8\u207d\u2071\u207e)^\u0394\u03b2",
        13,
    ),
    (
        " 7:      Normalize: W\u207d\u2071\u207e = w\u207d\u2071\u207e / \u2211\u2096 w\u207d\u1d4f\u207e",
        13,
    ),
    (" 8:      Update evidence: p\u0302(D) \u00d7= (1/N) \u2211\u1d62 w\u207d\u2071\u207e", 13),
    (" 9:      Resample N particles ~ Multinomial(W)", 13),
    ("10:      Compute \u03a3\u2c7c = \u03b2\u00b2 \u00b7 Cov_w(\u03b8)  (adapted proposal)", 13),
    ("11:      for each particle i = 1,...,N do", 13),
    ("12:          Metropolis-Hastings step with N(\u03b8\u207d\u2071\u207e, \u03a3\u2c7c)", 13),
    ("13:      end for", 13),
    ("14:      j = j + 1", 13),
    ("15:  end while", 13),
    ("16:  Return {\u03b8\u207d\u2071\u207e}, p\u0302(D)", 13),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(8.0), Inches(5.5), algo, font="Consolas")

_card(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(5.7))
_txt(
    s,
    Inches(9.4),
    Inches(1.4),
    Inches(3.4),
    Inches(0.35),
    "Key Equations / 主要な数式",
    sz=15,
    bold=True,
    color=TITLE_BLUE,
)
_txt(
    s,
    Inches(9.4),
    Inches(1.75),
    Inches(3.4),
    Inches(0.25),
    "Step 6-7: Importance Weights",
    sz=12,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_weights.png"), Inches(9.3), Inches(2.0), width=Inches(3.6))
_txt(
    s,
    Inches(9.4),
    Inches(2.85),
    Inches(3.4),
    Inches(0.25),
    "Step 10: Proposal Covariance",
    sz=12,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_covariance.png"), Inches(9.3), Inches(3.1), width=Inches(3.6))
_txt(
    s,
    Inches(9.4),
    Inches(3.8),
    Inches(3.4),
    Inches(0.25),
    "Step 12: MH Acceptance",
    sz=12,
    bold=True,
    color=DARK_TEXT,
)
_img(s, os.path.join(EQIMG, "eq_mh_accept.png"), Inches(9.3), Inches(4.05), width=Inches(3.6))
lines = [
    ("", 4),
    ("Step 8: Evidence update", 12, True, DARK_TEXT),
    ("accumulated as by-product", 11, False, MID_GRAY),
    ("Step 9: Multinomial", 12, True, DARK_TEXT),
    ("resampling eliminates", 11, False, MID_GRAY),
    ("low-weight particles", 11, False, MID_GRAY),
]
_multi(s, Inches(9.4), Inches(4.9), Inches(3.4), Inches(1.8), lines)
_footer(s)


# -- Model evidence slide --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Model Evidence & Bayes Factor", "モデルエビデンスとベイズ因子")

_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(3.0), TITLE_BLUE)
_txt(
    s,
    Inches(0.7),
    Inches(1.4),
    Inches(5.8),
    Inches(0.4),
    "Model Evidence (TMCMC by-product)",
    sz=20,
    bold=True,
    color=TITLE_BLUE,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(0.7),
    Inches(1.8),
    Inches(5.8),
    Inches(0.3),
    "モデルエビデンス (TMCMC副産物)",
    sz=13,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_img(s, os.path.join(EQIMG, "eq_evidence.png"), Inches(0.8), Inches(2.1), width=Inches(5.5))
lines = [
    ("Accumulated product of average weights", 13, False, MID_GRAY),
    ("across all tempering stages.", 13, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(3.4), Inches(5.8), Inches(0.6), lines, align=PP_ALIGN.CENTER)

_card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(3.0), ORANGE)
_txt(
    s,
    Inches(7.2),
    Inches(1.4),
    Inches(5.6),
    Inches(0.4),
    "Bayes Factor for Model Selection",
    sz=20,
    bold=True,
    color=ORANGE,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(7.2),
    Inches(1.8),
    Inches(5.6),
    Inches(0.3),
    "モデル選択のためのベイズ因子",
    sz=13,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_img(s, os.path.join(EQIMG, "eq_bayes_factor.png"), Inches(7.8), Inches(2.1), width=Inches(4.5))
lines = [
    ("BF > 10 : Strong evidence for M\u2081", 15),
    ("BF > 100 : Decisive evidence", 15),
    ("Compare: 15-param vs 20-param", 14, False, MID_GRAY),
]
_multi(s, Inches(7.2), Inches(3.3), Inches(5.6), Inches(0.8), lines, align=PP_ALIGN.CENTER)

_card(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.2))
lines = [
    ("Practical Significance / 実用的意義", 20, True, DARK_TEXT),
    ("", 6),
    (
        "\u2022 TMCMC produces model evidence automatically \u2014 no additional computation needed",
        16,
    ),
    ("  TMCMCはモデルエビデンスを自動的に算出 \u2014 追加計算不要", 13, False, MID_GRAY),
    (
        "\u2022 Enables principled comparison between constrained (15-param) and full (20-param) models",
        16,
    ),
    ("  制約付き(15パラメータ)と全(20パラメータ)モデルの原理的比較が可能", 13, False, MID_GRAY),
    ("\u2022 Quantitative justification for biological parameter locking", 16),
    ("  生物学的パラメータロックの定量的正当化", 13, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(4.8), Inches(12), Inches(2.0), lines)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 6: SEQUENTIAL ESTIMATION
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u2465 4-Stage Sequential Estimation", "4段階逐次推定アルゴリズム")

s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Algorithm 2: Sequential Estimation", "逐次推定アルゴリズム")

_card(s, Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.7), TITLE_BLUE)
algo2 = [
    ("Algorithm 2: 4-Stage Sequential Parameter Estimation", 16, True, TITLE_BLUE),
    ("\u2500" * 55, 10, False, LIGHT_GRAY),
    ("Input:  4 datasets {D_k}, lock rules {L_k} for k=1,...,4", 13),
    ("Output: Condition-specific posteriors {p(\u03b8|D_k)}", 13),
    ("", 6),
    ("Stage 1: Commensal Static  [N_locked = 9]", 14, True, GREEN),
    ("  \u2022 Prior: Uniform[\u2212100, 100] for interactions", 13),
    ("  \u2022 Lock: standard 5 + pathogen interactions", 13),
    ("  \u2022 Run TMCMC(D\u2081, L\u2081) \u2192 \u03b8\u2081 posterior", 13),
    ("", 4),
    ("Stage 2: Dysbiotic Static  [N_locked = 5]", 14, True, ORANGE),
    ("  \u2022 Prior: Uniform[\u2212100, 100]", 13),
    ("  \u2022 Lock: standard 5 biological locks", 13),
    ("  \u2022 Run TMCMC(D\u2082, L\u2082) \u2192 \u03b8\u2082 posterior", 13),
    ("", 4),
    ("Stage 3: Commensal HOBIC  [N_locked = 8]", 14, True, TITLE_BLUE),
    ("  \u2022 Strict locks with S.o growth allowed", 13),
    ("  \u2022 Run TMCMC(D\u2083, L\u2083) \u2192 \u03b8\u2083 posterior", 13),
    ("", 4),
    ("Stage 4: Dysbiotic HOBIC (Surge)  [N_locked = 0]", 14, True, ACCENT_RED),
    ("  \u2605 UNLOCK ALL \u2014 Discovery Mode", 13, True, ACCENT_RED),
    ("  \u2022 All 20 parameters free", 13),
    ("  \u2022 Run TMCMC(D\u2084, \u2205) \u2192 \u03b8\u2084 posterior", 13),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(7.5), Inches(5.5), algo2, font="Consolas")

_card(s, Inches(8.7), Inches(1.3), Inches(4.3), Inches(5.7))
lines = [
    ("Design Rationale", 16, True, TITLE_BLUE),
    ("設計の根拠", 12, False, MID_GRAY),
    ("", 8),
    ("\u2460 Start with most", 14, True, DARK_TEXT),
    ("   constrained condition", 14),
    ("   (CS: 9 locks = easiest)", 13, False, MID_GRAY),
    ("", 6),
    ("\u2461 Progressively relax", 14, True, DARK_TEXT),
    ("   constraints as model", 14),
    ("   complexity increases", 14),
    ("", 6),
    ("\u2462 Each stage provides", 14, True, DARK_TEXT),
    ("   independent posterior", 14),
    ("   for that condition", 14),
    ("", 6),
    ("\u2463 Final stage: Discovery", 14, True, ACCENT_RED),
    ("   Mode captures complex", 14),
    ("   Surge dynamics with", 14),
    ("   all parameters free", 14),
    ("", 8),
    ("N_particles = 1000", 14, True, DARK_TEXT),
    ("per condition", 14, False, MID_GRAY),
]
_multi(s, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.4), lines)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 7: EXPERIMENTAL CONDITIONS OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════

# -- 4 conditions overview with target data --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Experimental Target Data", "実験ターゲットデータ — 4条件")

positions = [
    ("Commensal_Static", Inches(0.2), Inches(1.2)),
    ("Dysbiotic_Static", Inches(6.7), Inches(1.2)),
    ("Commensal_HOBIC", Inches(0.2), Inches(4.2)),
    ("Dysbiotic_HOBIC", Inches(6.7), Inches(4.2)),
]
for key, x, y in positions:
    c = CONDS[key]
    tgt = os.path.join(EXPFIG, f"Target_{key}.png")
    _txt(
        s,
        x + Inches(0.1),
        y,
        Inches(4),
        Inches(0.3),
        c["label"],
        sz=13,
        bold=True,
        color=c["color"],
    )
    _txt(
        s,
        x + Inches(4.1),
        y,
        Inches(2),
        Inches(0.3),
        f"N_locked={c['locked']}",
        sz=12,
        bold=True,
        color=ACCENT_RED if c["locked"] == 0 else GREEN,
        align=PP_ALIGN.RIGHT,
    )
    _img(s, tgt, x + Inches(0.1), y + Inches(0.3), width=Inches(6.1), height=Inches(2.5))
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# RESULTS: 4 slides per condition (Overview, Heatmap+Violin, Diagnostics, Extra)
# ═══════════════════════════════════════════════════════════════════════════════
section_nums = {
    "Commensal_Static": "\u2466",
    "Dysbiotic_Static": "\u2467",
    "Commensal_HOBIC": "\u2468",
    "Dysbiotic_HOBIC": "\u2469",
}

for cond_key, cond in CONDS.items():
    d = cond["dir"]
    sec = section_nums[cond_key]

    # ── Section header ──
    _section(f"{sec} Results: {cond['label']}", cond["label_ja"])

    # ── Slide A: Per-species posterior fit (main result) ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(
        s,
        f"{cond['label']}  \u2014  Posterior Predictive Fit",
        f"{cond['label_ja']}  \u2014  事後予測適合",
    )

    # Description card (left)
    _card(s, Inches(0.4), Inches(1.3), Inches(4.2), Inches(5.5), cond["color"])
    lines = [
        ("Condition", 18, True, cond["color"]),
        (cond["desc_en"], 13),
        ("", 4),
        (cond["desc_ja"], 12, False, MID_GRAY),
        ("", 8),
        (f"N_locked = {cond['locked']}", 22, True, ACCENT_RED if cond["locked"] == 0 else GREEN),
        ("", 8),
        ("Key Finding", 18, True, cond["color"]),
        (cond["finding_en"], 13),
        ("", 4),
        (cond["finding_ja"], 12, False, MID_GRAY),
    ]
    _multi(s, Inches(0.7), Inches(1.5), Inches(3.8), Inches(5.2), lines)

    # Per-species panel (right, large)
    fig = os.path.join(d, "Fig_A02_per_species_panel.png")
    _img(s, fig, Inches(4.8), Inches(1.2), width=Inches(8.2))
    _txt(
        s,
        Inches(4.8),
        Inches(6.8),
        Inches(8.2),
        Inches(0.3),
        "Per-species posterior predictive fit with 95% credible interval / 菌種別事後予測適合 (95%信用区間)",
        sz=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _footer(s)

    # ── Slide B: Interaction heatmap + Violin plot ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(
        s,
        f"{cond['label']}  \u2014  Interaction Matrix & Parameter Uncertainty",
        f"{cond['label_ja']}  \u2014  相互作用行列 & パラメータ不確実性",
    )

    heatmap = os.path.join(d, "Fig_A01_interaction_matrix_heatmap.png")
    violin = os.path.join(d, "Fig_A05_parameter_violins.png")

    _img(s, heatmap, Inches(0.3), Inches(1.2), width=Inches(6.2))
    _txt(
        s,
        Inches(0.3),
        Inches(6.5),
        Inches(6.2),
        Inches(0.4),
        "Estimated interaction matrix A (MAP) / 推定相互作用行列",
        sz=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, violin, Inches(6.7), Inches(1.2), width=Inches(6.3))
    _txt(
        s,
        Inches(6.7),
        Inches(6.5),
        Inches(6.3),
        Inches(0.4),
        "Parameter posterior distributions (violin) / パラメータ事後分布 (バイオリン)",
        sz=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _footer(s)

    # ── Slide C: MAP vs Mean + Residuals ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(
        s,
        f"{cond['label']}  \u2014  MAP vs Mean & Residuals",
        f"{cond['label_ja']}  \u2014  MAP vs 平均推定 & 残差分析",
    )

    map_mean = os.path.join(d, "Fig_A09_MAP_vs_Mean_comparison.png")
    residual = os.path.join(d, "figures", f"residuals_{cond_key}_Residuals.png")

    _img(s, map_mean, Inches(0.3), Inches(1.2), width=Inches(6.2))
    _txt(
        s,
        Inches(0.3),
        Inches(6.5),
        Inches(6.2),
        Inches(0.4),
        "MAP vs Mean parameter comparison / MAP vs 平均パラメータ比較",
        sz=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, residual, Inches(6.7), Inches(1.2), width=Inches(6.3))
    _txt(
        s,
        Inches(6.7),
        Inches(6.5),
        Inches(6.3),
        Inches(0.4),
        "Residual analysis / 残差分析",
        sz=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _footer(s)

    # ── Slide D: Convergence + Composition + Correlation ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(
        s,
        f"{cond['label']}  \u2014  Convergence & Diagnostics",
        f"{cond['label_ja']}  \u2014  収束診断",
    )

    beta_fig = os.path.join(d, "figures", f"Fig01_TMCMC_beta_schedule_{cond_key}.png")
    comp_fig = os.path.join(d, "Fig_A04_species_composition.png")
    corr_fig = os.path.join(d, "Fig_A06_correlation_matrix.png")
    conv_fig = os.path.join(d, "Fig_A10_convergence_dashboard.png")

    # Top row: 3 figures
    _img(s, beta_fig, Inches(0.3), Inches(1.2), width=Inches(4.1), height=Inches(2.7))
    _txt(
        s,
        Inches(0.3),
        Inches(3.85),
        Inches(4.1),
        Inches(0.3),
        "TMCMC \u03b2 schedule",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, comp_fig, Inches(4.6), Inches(1.2), width=Inches(4.1), height=Inches(2.7))
    _txt(
        s,
        Inches(4.6),
        Inches(3.85),
        Inches(4.1),
        Inches(0.3),
        "Species composition / 菌種構成",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, corr_fig, Inches(8.9), Inches(1.2), width=Inches(4.1), height=Inches(2.7))
    _txt(
        s,
        Inches(8.9),
        Inches(3.85),
        Inches(4.1),
        Inches(0.3),
        "Parameter correlation / パラメータ相関",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    # Bottom: convergence dashboard
    _img(s, conv_fig, Inches(0.3), Inches(4.4), width=Inches(12.7), height=Inches(2.3))
    _txt(
        s,
        Inches(0.3),
        Inches(6.75),
        Inches(12.7),
        Inches(0.3),
        "Convergence dashboard / 収束ダッシュボード",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _footer(s)

    # ── Slide E: Extra diagnostics (state decomposition, log-L, posterior predictive check) ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(
        s,
        f"{cond['label']}  \u2014  Additional Diagnostics",
        f"{cond['label_ja']}  \u2014  追加診断",
    )

    state_fig = os.path.join(d, "Fig_A03_state_decomposition.png")
    logl_fig = os.path.join(d, "Fig_A07_loglikelihood_landscape.png")
    ppc_fig = os.path.join(d, "Fig_A08_posterior_predictive_check.png")

    _img(s, state_fig, Inches(0.3), Inches(1.2), width=Inches(4.1))
    _txt(
        s,
        Inches(0.3),
        Inches(6.5),
        Inches(4.1),
        Inches(0.3),
        "State decomposition (\u03c6, \u03c8) / 状態分解",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, logl_fig, Inches(4.6), Inches(1.2), width=Inches(4.1))
    _txt(
        s,
        Inches(4.6),
        Inches(6.5),
        Inches(4.1),
        Inches(0.3),
        "Log-likelihood landscape / 対数尤度",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )

    _img(s, ppc_fig, Inches(8.9), Inches(1.2), width=Inches(4.1))
    _txt(
        s,
        Inches(8.9),
        Inches(6.5),
        Inches(4.1),
        Inches(0.3),
        "Posterior predictive check / 事後予測チェック",
        sz=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 11: COMPARATIVE ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u246A Comparative Analysis", "比較分析")

# -- Heatmaps side by side --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Interaction Matrix Comparison (All 4 Conditions)", "相互作用行列の比較 \u2014 全4条件")

hm_pos = [
    ("Commensal_Static", Inches(0.15), "CS"),
    ("Dysbiotic_Static", Inches(3.35), "DS"),
    ("Commensal_HOBIC", Inches(6.55), "CH"),
    ("Dysbiotic_HOBIC", Inches(9.75), "DH"),
]
for key, x, short in hm_pos:
    c = CONDS[key]
    hm = os.path.join(c["dir"], "Fig_A01_interaction_matrix_heatmap.png")
    _img(s, hm, x, Inches(1.3), width=Inches(3.1))
    _txt(
        s,
        x,
        Inches(5.5),
        Inches(3.1),
        Inches(0.5),
        f"{c['label']}  (N_locked={c['locked']})",
        sz=11,
        color=c["color"],
        align=PP_ALIGN.CENTER,
        bold=True,
    )

lines = [
    (
        "\u2022 Commensal: competitive (blue) blocks dominate  \u2022 Dysbiotic: cooperative (red) blocks emerge around P.g",
        14,
        False,
        DARK_TEXT,
    ),
    (
        "\u2022 健康: 競合的(青)ブロックが支配  \u2022 疾患: P.g周辺に協力的(赤)ブロックが出現",
        12,
        False,
        MID_GRAY,
    ),
]
_multi(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.8), lines)
_footer(s)

# -- Per-species fit comparison --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(
    s, "Posterior Predictive Fit Comparison (All 4 Conditions)", "事後予測適合の比較 \u2014 全4条件"
)

pp_pos = [
    ("Commensal_Static", Inches(0.15), Inches(1.2)),
    ("Dysbiotic_Static", Inches(6.75), Inches(1.2)),
    ("Commensal_HOBIC", Inches(0.15), Inches(4.2)),
    ("Dysbiotic_HOBIC", Inches(6.75), Inches(4.2)),
]
for key, x, y in pp_pos:
    c = CONDS[key]
    fig = os.path.join(c["dir"], "Fig_A02_per_species_panel.png")
    _txt(
        s,
        x + Inches(0.1),
        y,
        Inches(5),
        Inches(0.3),
        c["label"],
        sz=12,
        bold=True,
        color=c["color"],
    )
    _img(s, fig, x, y + Inches(0.3), width=Inches(6.4), height=Inches(2.55))
_footer(s)

# -- Violin comparison --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(
    s,
    "Parameter Uncertainty Comparison (All 4 Conditions)",
    "パラメータ不確実性の比較 \u2014 全4条件",
)

vl_pos = [
    ("Commensal_Static", Inches(0.15), Inches(1.2)),
    ("Dysbiotic_Static", Inches(6.75), Inches(1.2)),
    ("Commensal_HOBIC", Inches(0.15), Inches(4.2)),
    ("Dysbiotic_HOBIC", Inches(6.75), Inches(4.2)),
]
for key, x, y in vl_pos:
    c = CONDS[key]
    fig = os.path.join(c["dir"], "Fig_A05_parameter_violins.png")
    _txt(
        s,
        x + Inches(0.1),
        y,
        Inches(5),
        Inches(0.3),
        c["label"],
        sz=12,
        bold=True,
        color=c["color"],
    )
    _img(s, fig, x, y + Inches(0.3), width=Inches(6.4), height=Inches(2.55))
_footer(s)

# -- Commensal vs Dysbiotic analysis --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Commensal vs. Dysbiotic  \u2014  Key Differences", "健康 vs 疾患 \u2014 主な違い")

_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.5), GREEN)
lines = [
    ("Commensal (Healthy)", 24, True, GREEN),
    ("健康状態", 14, False, MID_GRAY),
    ("", 8),
    ("\u2022 Early colonizers (S.o, A.n) dominate", 16),
    ("  初期定着菌 (S.o, A.n) が優位", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Pathogen (P.g) interactions suppressed", 16),
    ("  病原菌 (P.g) の相互作用が抑制", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Interaction matrix: competitive (blue)", 16),
    ("  相互作用行列: 競合的 (青) ブロックが支配", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Narrow posteriors = high identifiability", 16),
    ("  狭い事後分布 = 高い識別性", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Strict locks reinforce biological knowledge", 16),
    ("  厳格なロックが生物学的知識を強化", 13, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.2), lines)

_card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(5.5), ACCENT_RED)
lines = [
    ("Dysbiotic (Disease)", 24, True, ACCENT_RED),
    ("疾患状態", 14, False, MID_GRAY),
    ("", 8),
    ("\u2022 P. gingivalis interactions become positive", 16),
    ("  P.g の相互作用が正 (協力的) に転換", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Veillonella\u2013P.g cross-feeding emerges", 16),
    ("  Vei\u2013P.g クロスフィーディングが出現", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Full dysbiosis requires flow (HOBIC)", 16),
    ("  完全な疾患にはHOBIC (流れ) が必要", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 'Surge' in HOBIC: explosive P.g growth", 16, True, ACCENT_RED),
    ("  HOBIC での 'Surge': P.g の爆発的増殖", 13, False, MID_GRAY),
    ("", 4),
    ("\u2022 Discovery Mode essential to capture dynamics", 16),
    ("  Discovery Mode が動態の捕捉に不可欠", 13, False, MID_GRAY),
]
_multi(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(5.2), lines)
_footer(s)

# -- Static vs HOBIC --
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Static vs. HOBIC  \u2014  Effect of Saliva Flow", "静的 vs 流動 \u2014 唾液流の影響")

_card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.5))
lines = [
    ("Static Culture", 22, True, DARK_TEXT),
    ("静的培養", 14, False, MID_GRAY),
    ("", 8),
    ("\u2022 Nutrient-limited environment", 16),
    ("  栄養制限環境", 13, False, MID_GRAY),
    ("\u2022 Stable but low biomass steady states", 16),
    ("  安定だが低バイオマスの定常状態", 13, False, MID_GRAY),
    ("\u2022 Metabolite accumulation limits pathogens", 16),
    ("  代謝産物蓄積が病原菌を制限", 13, False, MID_GRAY),
    ("\u2022 Narrower posterior distributions", 16),
    ("  より狭い事後分布", 13, False, MID_GRAY),
    ("\u2022 Less complex dynamics \u2192 easier estimation", 16),
    ("  より単純な動態 \u2192 容易な推定", 13, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.2), lines)

_card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(5.5))
lines = [
    ("HOBIC (Flow)", 22, True, TITLE_BLUE),
    ("流動培養 (HOBIC)", 14, False, MID_GRAY),
    ("", 8),
    ("\u2022 Simulates oral saliva flow", 16),
    ("  口腔の唾液流を模倣", 13, False, MID_GRAY),
    ("\u2022 More dynamic steady states", 16),
    ("  より動的な定常状態", 13, False, MID_GRAY),
    ("\u2022 Higher variance in decay rate posteriors", 16),
    ("  減衰率事後分布の分散が増大", 13, False, MID_GRAY),
    ("\u2022 Enables 'Blue Bloom' (Commensal)", 16, False, GREEN),
    ("  'Blue Bloom' を可能にする (健康)", 13, False, MID_GRAY),
    ("\u2022 Enables 'Surge' (Dysbiotic)", 16, False, ACCENT_RED),
    ("  'Surge' を可能にする (疾患)", 13, False, MID_GRAY),
]
_multi(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(5.2), lines)
_footer(s)

# -- Cross-condition comparison figures --
pub_map = os.path.join(DOCS, "pub_map_fit_comparison.png")
pub_hm = os.path.join(DOCS, "pub_interaction_heatmap.png")
pub_fit = os.path.join(DOCS, "pub_fit_comparison.png")

if any(os.path.exists(f) for f in [pub_map, pub_hm, pub_fit]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "Publication-Quality Comparison Figures", "論文品質の比較図")

    figs = [
        (pub_map, "MAP fit comparison"),
        (pub_hm, "Interaction heatmap comparison"),
        (pub_fit, "Fit comparison"),
    ]
    x_pos = Inches(0.3)
    for path, caption in figs:
        if os.path.exists(path):
            _img(s, path, x_pos, Inches(1.2), width=Inches(4.1))
            _txt(
                s,
                x_pos,
                Inches(6.5),
                Inches(4.1),
                Inches(0.3),
                caption,
                sz=11,
                color=MID_GRAY,
                align=PP_ALIGN.CENTER,
            )
            x_pos += Inches(4.3)
    _footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 12: CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
_section("\u246B Conclusion & Future Work", "結論・今後の課題")

s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s)
_header(s, "Conclusion / 結論")

_card(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(3.0), TITLE_BLUE)
lines = [
    ("Key Contributions / 主な貢献", 22, True, TITLE_BLUE),
    ("", 6),
    (
        "\u2460  Biologically-constrained parameter reduction: 30 \u2192 20 \u2192 15 free parameters",
        17,
        False,
        DARK_TEXT,
    ),
    (
        "    生物学的制約に基づくパラメータ削減: 30 \u2192 20 \u2192 15 自由パラメータ",
        13,
        False,
        MID_GRAY,
    ),
    ("", 4),
    (
        "\u2461  TMCMC successfully applied to all 4 experimental conditions with 1000 particles",
        17,
        False,
        DARK_TEXT,
    ),
    ("    TMCMCを1000粒子で4つの全実験条件に適用成功", 13, False, MID_GRAY),
    ("", 4),
    (
        "\u2462  Unlock All (Discovery Mode) essential for Dysbiotic HOBIC 'Surge' detection",
        17,
        False,
        DARK_TEXT,
    ),
    (
        "    全ロック解除 (Discovery Mode) がDysbiotic HOBIC 'Surge' の検出に不可欠",
        13,
        False,
        MID_GRAY,
    ),
    ("", 4),
    (
        "\u2463  Inferred interaction matrices provide quantitative health \u2192 disease transition map",
        17,
        False,
        DARK_TEXT,
    ),
    ("    推定相互作用行列が健康 \u2192 疾患遷移の定量的マップを提供", 13, False, MID_GRAY),
    ("", 4),
    (
        "\u2464  4-stage sequential estimation enables condition-specific Bayesian inference",
        17,
        False,
        DARK_TEXT,
    ),
    ("    4段階逐次推定が条件特異的ベイズ推論を実現", 13, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(1.5), Inches(12), Inches(2.8), lines)

_card(s, Inches(0.4), Inches(4.6), Inches(6.0), Inches(2.2), ORANGE)
lines = [
    ("Future Work / 今後の課題", 20, True, ORANGE),
    ("", 6),
    ("\u2022 Increase particle count (N > 5000)", 15),
    ("  粒子数の増加 (N > 5000)", 12, False, MID_GRAY),
    ("\u2022 Time-varying interaction parameters", 15),
    ("  時間変化する相互作用パラメータ", 12, False, MID_GRAY),
    ("\u2022 In silico therapeutic testing", 15),
    ("  治療介入のin silicoテスト", 12, False, MID_GRAY),
]
_multi(s, Inches(0.7), Inches(4.8), Inches(5.6), Inches(2.0), lines)

_card(s, Inches(6.7), Inches(4.6), Inches(6.2), Inches(2.2), GREEN)
lines = [
    ("Clinical Impact / 臨床的インパクト", 20, True, GREEN),
    ("", 6),
    ("\u2022 Quantitative dysbiosis biomarkers", 15),
    ("  定量的ディスバイオシスバイオマーカー", 12, False, MID_GRAY),
    ("\u2022 Personalized peri-implantitis risk", 15),
    ("  個別化インプラント周囲炎リスク", 12, False, MID_GRAY),
    ("\u2022 Framework for multi-species analysis", 15),
    ("  多菌種解析の枠組み", 12, False, MID_GRAY),
]
_multi(s, Inches(7.0), Inches(4.8), Inches(5.8), Inches(2.0), lines)
_footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg(s, NAVY)
_rect(s, Inches(0), Inches(0), SW, Inches(0.08), GOLD)
_rect(s, Inches(0), SH - Inches(0.08), SW, Inches(0.08), GOLD)
_txt(
    s,
    Inches(1),
    Inches(2.0),
    Inches(11.3),
    Inches(1.5),
    "Thank You",
    sz=54,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(1),
    Inches(3.5),
    Inches(11.3),
    Inches(0.8),
    "ご清聴ありがとうございます",
    sz=28,
    color=GOLD,
    align=PP_ALIGN.CENTER,
)
_rect(s, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.03), TITLE_BLUE)
_txt(
    s,
    Inches(1),
    Inches(4.8),
    Inches(11.3),
    Inches(1.0),
    "Keisuke Nishioka\nInstitut f\u00fcr Kontinuumsmechanik (IKM)\nLeibniz Universit\u00e4t Hannover",
    sz=18,
    color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER,
)
_txt(
    s,
    Inches(1),
    Inches(6.3),
    Inches(11.3),
    Inches(0.4),
    "February 2026",
    sz=14,
    color=MID_GRAY,
    align=PP_ALIGN.CENTER,
)
_footer(s)


# ── Save ─────────────────────────────────────────────────────────────────────
output = os.path.join(DOCS, "nishioka_biofilm_tmcmc_presentation.pptx")
prs.save(output)
print(f"\nPresentation saved: {output}")
print(f"Total slides: {len(prs.slides)}")
print(f"Slide counter: {slide_counter[0]}")
