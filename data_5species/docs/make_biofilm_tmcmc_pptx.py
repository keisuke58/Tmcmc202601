#!/usr/bin/env python3
"""
Professional conference presentation: 5-species biofilm FOM / ROM / TMCMC.

All equations are rendered as LaTeX images via matplotlib (usetex).
Generates both Japanese and English 16:9 PPTX files with speaker notes.
"""

from __future__ import annotations

import hashlib
import os

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from matplotlib import rc as mplrc  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

# ── paths ──────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
EQ_DIR = os.path.join(SCRIPT_DIR, "_eq_img")
os.makedirs(EQ_DIR, exist_ok=True)

# ── slide size (16 : 9) ───────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── layout constants ──────────────────────────────────────────
MARGIN_L = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN_L
TITLE_BAR_H = Inches(1.05)
CONTENT_Y0 = Inches(1.40)
FOOTER_Y = Inches(7.05)

# ── colours ───────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x22, 0x22, 0x22)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x00, 0x72, 0xCE)
SUBTLE = RGBColor(0xAA, 0xBB, 0xDD)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)

# ── fonts ─────────────────────────────────────────────────────
FONT = "Calibri"

# ── LaTeX rendering ───────────────────────────────────────────
EQ_DPI = 300
EQ_FONTSIZE = 24

mplrc("text", usetex=True)
mplrc("font", family="serif")
mplrc(
    "text.latex",
    preamble=r"\usepackage{amsmath}\usepackage{amssymb}\usepackage{bm}\usepackage{xcolor}",
)


# ═══════════════════════════════════════════════════════════════
#  LaTeX -> PNG helpers
# ═══════════════════════════════════════════════════════════════
def _render_latex(
    tex: str,
    fontsize: int = EQ_FONTSIZE,
    dpi: int = EQ_DPI,
    color: str = "black",
) -> str:
    """Render *tex* to a transparent PNG; return the cached file path."""
    tag = hashlib.md5(f"{tex}|{fontsize}|{dpi}|{color}".encode()).hexdigest()[:12]
    path = os.path.join(EQ_DIR, f"eq_{tag}.png")
    if os.path.exists(path):
        return path
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(
        0,
        0,
        tex,
        fontsize=fontsize,
        color=color,
        verticalalignment="bottom",
        horizontalalignment="left",
    )
    fig.savefig(path, dpi=dpi, bbox_inches="tight", pad_inches=0.04, transparent=True)
    plt.close(fig)
    return path


def _img_emu(path: str, scale: float = 1.0) -> tuple[int, int]:
    """Return (width, height) in EMU for a rendered equation image."""
    with Image.open(path) as img:
        w_px, h_px = img.size
    return (
        int(w_px / EQ_DPI * scale * 914400),
        int(h_px / EQ_DPI * scale * 914400),
    )


# ═══════════════════════════════════════════════════════════════
#  Slide-building helpers
# ═══════════════════════════════════════════════════════════════
_slide_counter: int = 0


def _new_slide(prs: Presentation, title: str):
    """Create a blank slide with a navy title bar + accent line."""
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # title-bar background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # accent line below bar
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, TITLE_BAR_H, SLIDE_W, Inches(0.04))
    acc.fill.solid()
    acc.fill.fore_color.rgb = ACCENT
    acc.line.fill.background()

    # title text
    tb = slide.shapes.add_textbox(MARGIN_L, Inches(0.18), CONTENT_W, Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.color.rgb = WHITE
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = FONT

    # slide number (footer right)
    sn = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), FOOTER_Y, Inches(0.6), Inches(0.3))
    sp = sn.text_frame.paragraphs[0]
    sp.text = str(_slide_counter)
    sp.alignment = PP_ALIGN.RIGHT
    sp.font.size = Pt(12)
    sp.font.color.rgb = MID_GRAY
    sp.font.name = FONT

    return slide


# ── content primitives ────────────────────────────────────────
def _gap(n: float = 0.15) -> int:
    return Inches(n)


def _add_text(
    slide,
    text: str,
    top: int,
    *,
    left: int | None = None,
    width: int | None = None,
    size: int = Pt(20),
    bold: bool = False,
    color: RGBColor = NEAR_BLACK,
    align: int = PP_ALIGN.LEFT,
) -> int:
    """Add a text box. Returns the estimated bottom-y (EMU)."""
    if left is None:
        left = MARGIN_L
    if width is None:
        width = CONTENT_W
    lines = text.split("\n")
    height = int(size * 1.8 * len(lines))
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT
        p.alignment = align
    return top + height


def _add_bullets(
    slide,
    items: list[str],
    top: int,
    *,
    size: int = Pt(20),
    color: RGBColor = NEAR_BLACK,
) -> int:
    """Add a bulleted list. Returns estimated bottom-y."""
    height = int(size * 2.0 * len(items))
    tb = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(4)
    return top + height


def _add_eq(
    slide,
    tex: str,
    top: int,
    *,
    fontsize: int = EQ_FONTSIZE,
    scale: float = 1.0,
    center: bool = True,
    max_w: int | None = None,
) -> int:
    """Render LaTeX, insert as picture. Returns bottom-y."""
    path = _render_latex(tex, fontsize=fontsize)
    w, h = _img_emu(path, scale=scale)
    limit = max_w if max_w else int(CONTENT_W)
    if w > limit:
        ratio = limit / w
        w = limit
        h = int(h * ratio)
    if center:
        left = MARGIN_L + (int(CONTENT_W) - w) // 2
    else:
        left = MARGIN_L + Inches(0.3)
    slide.shapes.add_picture(path, left, top, w, h)
    return top + h


def _add_eq_at(
    slide,
    tex: str,
    left: int,
    top: int,
    *,
    fontsize: int = EQ_FONTSIZE,
    scale: float = 1.0,
    max_w: int | None = None,
) -> tuple[int, int]:
    """Render LaTeX at explicit (left, top). Returns (w, h) in EMU."""
    path = _render_latex(tex, fontsize=fontsize)
    w, h = _img_emu(path, scale=scale)
    if max_w and w > max_w:
        ratio = max_w / w
        w = max_w
        h = int(h * ratio)
    slide.shapes.add_picture(path, left, top, w, h)
    return w, h


def _add_symbol_list(
    slide,
    items: list[tuple[str, str]],
    top: int,
    eq_fontsize: int = 20,
    desc_size: int = Pt(18),
) -> int:
    """Add (LaTeX symbol, description) pairs in two-column layout."""
    EQ_COL_W = Inches(2.8)
    y = top
    for tex, desc in items:
        path = _render_latex(tex, fontsize=eq_fontsize)
        w, h = _img_emu(path)
        row_h = max(h, int(desc_size * 2.0))
        # equation image (left)
        slide.shapes.add_picture(
            path,
            MARGIN_L + Inches(0.3),
            y + (row_h - h) // 2,
            w,
            h,
        )
        # description text (right)
        desc_left = MARGIN_L + EQ_COL_W
        desc_w = int(CONTENT_W) - int(EQ_COL_W)
        tb = slide.shapes.add_textbox(desc_left, y, desc_w, row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = desc_size
        p.font.color.rgb = NEAR_BLACK
        p.font.name = FONT
        y += row_h + _gap(0.03)
    return y


def _add_matrix_pair(
    slide,
    a_tex: str,
    a_fs: int,
    b_tex: str,
    b_fs: int,
    top: int,
) -> int:
    """Render two equations side by side. Returns bottom-y."""
    a_path = _render_latex(a_tex, fontsize=a_fs)
    aw, ah = _img_emu(a_path)
    max_aw = int(CONTENT_W * 0.55)
    if aw > max_aw:
        ratio = max_aw / aw
        aw = max_aw
        ah = int(ah * ratio)
    slide.shapes.add_picture(a_path, MARGIN_L + Inches(0.3), top, aw, ah)

    b_path = _render_latex(b_tex, fontsize=b_fs)
    bw, bh = _img_emu(b_path)
    b_top = top + (ah - bh) // 2 if ah > bh else top
    slide.shapes.add_picture(b_path, MARGIN_L + int(CONTENT_W * 0.6), b_top, bw, bh)
    return top + max(ah, bh)


# ═══════════════════════════════════════════════════════════════
#  Generic slide renderer
# ═══════════════════════════════════════════════════════════════
# Element tuple formats:
#   ("h",  text)                         Header (bold, accent)
#   ("t",  text [, pt_size])             Text
#   ("tb", text [, pt_size])             Bold text
#   ("tc", text, RGBColor [, pt_size, bold])  Coloured text
#   ("b",  [items] [, pt_size])          Bullet list
#   ("e",  latex [, fontsize])           Centred equation
#   ("s",  [(latex, desc), ...] [, eq_fontsize])  Symbol list
#   ("g",  inches)                       Vertical gap
#   ("mp", a_tex, a_fs, b_tex, b_fs)    Matrix pair side-by-side


def _render_elements(slide, elements) -> int:
    """Render a list of content elements onto a slide. Returns final y."""
    y = CONTENT_Y0
    for el in elements:
        kind = el[0]
        if kind == "h":
            y = _add_text(slide, el[1], y, size=Pt(22), bold=True, color=ACCENT)
            y += _gap(0.05)
        elif kind == "t":
            sz = Pt(el[2]) if len(el) > 2 else Pt(20)
            y = _add_text(slide, el[1], y, size=sz)
        elif kind == "tb":
            sz = Pt(el[2]) if len(el) > 2 else Pt(20)
            y = _add_text(slide, el[1], y, size=sz, bold=True)
        elif kind == "tc":
            sz = Pt(el[3]) if len(el) > 3 else Pt(18)
            bl = el[4] if len(el) > 4 else True
            y = _add_text(slide, el[1], y, size=sz, bold=bl, color=el[2])
            y += _gap(0.05)
        elif kind == "b":
            sz = Pt(el[2]) if len(el) > 2 else Pt(20)
            y = _add_bullets(slide, el[1], y, size=sz)
        elif kind == "e":
            fs = el[2] if len(el) > 2 else EQ_FONTSIZE
            y = _add_eq(slide, el[1], y, fontsize=fs)
        elif kind == "s":
            fs = el[2] if len(el) > 2 else 20
            y = _add_symbol_list(slide, el[1], y, eq_fontsize=fs)
        elif kind == "g":
            y += _gap(el[1])
        elif kind == "mp":
            y = _add_matrix_pair(slide, el[1], el[2], el[3], el[4], y)
    return y


# ═══════════════════════════════════════════════════════════════
#  Title slide (special layout)
# ═══════════════════════════════════════════════════════════════
_TITLE = {
    "ja": {
        "lines": [
            "5種バイオフィルムモデルの",
            "FOM, ROM, TMCMC による",
            "パラメータ推定",
        ],
        "author": "西岡 佳祐",
        "affil": ("Institut f\u00fcr Kontinuumsmechanik, " "Leibniz Universit\u00e4t Hannover"),
        "date": "February 2026",
    },
    "en": {
        "lines": [
            "Parameter Estimation for a",
            "5-Species Biofilm Model via",
            "FOM, ROM, and TMCMC",
        ],
        "author": "Keisuke Nishioka",
        "affil": ("Institute of Continuum Mechanics, " "Leibniz University Hannover"),
        "date": "February 2026",
    },
}


def _slide_title(prs: Presentation, lang: str) -> None:
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    T = _TITLE[lang]

    # full navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    cx = (int(SLIDE_W) - Inches(11)) // 2
    cy = Inches(1.5)

    # title lines
    tb = slide.shapes.add_textbox(cx, cy, Inches(11), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(T["lines"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    # accent line
    lw = Inches(6)
    lx = (int(SLIDE_W) - lw) // 2
    ly = cy + Inches(2.8)
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, Inches(0.04))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT
    r.line.fill.background()

    # author block
    ab = slide.shapes.add_textbox(cx, ly + Inches(0.45), Inches(11), Inches(1.6))
    af = ab.text_frame
    af.word_wrap = True
    for i, (txt, sz, clr) in enumerate(
        [
            (T["author"], Pt(26), WHITE),
            (T["affil"], Pt(16), SUBTLE),
            (T["date"], Pt(16), SUBTLE),
        ]
    ):
        p = af.paragraphs[0] if i == 0 else af.add_paragraph()
        p.text = txt
        p.font.size = sz
        p.font.color.rgb = clr
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════
#  Shared equation strings (language-independent)
# ═══════════════════════════════════════════════════════════════
_EQ = dict(
    state_vec=(
        r"$\mathbf{g}(t)"
        r" = \bigl(\phi_1,\dots,\phi_5,\;\phi_0,\;"
        r"\psi_1,\dots,\psi_5,\;\gamma\bigr)^{\!\top}"
        r" \in \mathbb{R}^{12}$"
    ),
    vol_constraint=r"$\displaystyle\sum_{i=1}^{5}\phi_i(t) + \phi_0(t) = 1$",
    dissipative=(
        r"$\dfrac{d\mathbf{g}}{dt}"
        r" = -\,M(\mathbf{g})\;\nabla_{\!\mathbf{g}}\,F(\mathbf{g})"
        r" \;+\; S(\mathbf{g})$"
    ),
    implicit_euler=(
        r"$Q\!\bigl(\mathbf{g}^{\,n+1};\;"
        r"\mathbf{g}^{\,n},\,\boldsymbol{\theta}\bigr)"
        r" = \mathbf{0}$"
    ),
    res_phi_1=(r"$Q_{\phi_i}" r"\;=\;\dfrac{K_p\,(2 - 4\phi_i)}{(\phi_i - 1)^3\,\phi_i^3}$"),
    res_phi_2=(
        r"$\quad +\;\dfrac{1}{\eta_i}"
        r"\Bigl[\gamma"
        r" + \bigl(\eta_{\phi,i} + \eta_i\,\psi_i^2\bigr)\,\dot{\phi}_i"
        r" + \eta_i\,\phi_i\,\psi_i\,\dot{\psi}_i\Bigr]$"
    ),
    res_phi_3=(
        r"$\quad -\;\dfrac{c^{\!*}}{\eta_i}\,\psi_i"
        r"\,\bigl(A\,(\boldsymbol{\phi}"
        r"\odot\boldsymbol{\psi})\bigr)_{\!i}$"
    ),
    res_psi_1=(
        r"$Q_{\psi_i}"
        r"\;=\;-\dfrac{2K_p}{(\psi_i\!-\!1)^2\,\psi_i^3}"
        r" - \dfrac{2K_p}{(\psi_i\!-\!1)^3\,\psi_i^2}"
        r" \;+\; \dfrac{b_i\,\alpha^{\!*}}{\eta_i}\,\psi_i$"
    ),
    res_psi_2=(
        r"$\quad +\;\phi_i\,\psi_i\,\dot{\phi}_i"
        r" \;+\; \phi_i^2\,\dot{\psi}_i"
        r" \;-\; \dfrac{c^{\!*}}{\eta_i}\,\phi_i"
        r"\,\bigl(A\,(\boldsymbol{\phi}"
        r"\odot\boldsymbol{\psi})\bigr)_{\!i}$"
    ),
    res_gamma=(r"$Q_{\gamma}" r"\;=\;\displaystyle\sum_{i=1}^{5}\phi_i + \phi_0 - 1$"),
    theta_vec=(
        r"$\boldsymbol{\theta}"
        r" = \underbrace{(a_{11},a_{12},a_{22},b_1,b_2)}_{\mathrm{M1}}"
        r"\oplus \underbrace{(a_{33},a_{34},a_{44},b_3,b_4)}_{\mathrm{M2}}"
        r"\oplus \underbrace{(a_{13},a_{14},a_{23},a_{24})}_{\mathrm{M3}}"
        r"\oplus \underbrace{(a_{55},b_5)}_{\mathrm{M4}}"
        r"\oplus \underbrace{(a_{15},a_{25},a_{35},a_{45})}_{\mathrm{M5}}$"
    ),
    mat_A=(
        r"$A = \begin{pmatrix}"
        r"a_{11} & a_{12} & a_{13} & a_{14} & a_{15}\\"
        r"a_{12} & a_{22} & a_{23} & a_{24} & a_{25}\\"
        r"a_{13} & a_{23} & a_{33} & a_{34} & a_{35}\\"
        r"a_{14} & a_{24} & a_{34} & a_{44} & a_{45}\\"
        r"a_{15} & a_{25} & a_{35} & a_{45} & a_{55}"
        r"\end{pmatrix}$"
    ),
    mat_B=r"$B = \operatorname{diag}(b_1,\, b_2,\, b_3,\, b_4,\, b_5)$",
    locked_set=(
        r"$\mathcal{L} = \{6,\,12,\,13,\,16,\,17\}"
        r"\;\;\Longrightarrow\;\;"
        r"a_{34}=a_{23}=a_{24}=a_{15}=a_{25}=0$"
    ),
    n_free=(r"$n_{\mathrm{free}} = 20 - |\mathcal{L}| = 15$"),
    x_fn=r"$x_{\mathrm{FN}} = \phi_4\,\psi_4$",
    hill=(
        r"$h(x_{\mathrm{FN}})"
        r" = \dfrac{x_{\mathrm{FN}}^{\,n_{\mathrm{hill}}}}"
        r"{K_{\mathrm{hill}}^{\,n_{\mathrm{hill}}}"
        r" + x_{\mathrm{FN}}^{\,n_{\mathrm{hill}}}}$"
    ),
    a_eff=(r"$(A_{\mathrm{eff}})_{5j}" r" = h(x_{\mathrm{FN}})\;A_{5j}," r" \qquad j = 1,\dots,5$"),
    obs_model=(
        r"$y_{k,s}"
        r" = \hat{y}_{k,s}(\boldsymbol{\theta}) + \varepsilon_{k,s},"
        r"\qquad \varepsilon_{k,s}"
        r" \sim \mathcal{N}\!\bigl(0,\,\sigma_{\mathrm{obs}}^2\bigr)$"
    ),
    loglik=(
        r"$\log p(\mathbf{y}\mid\boldsymbol{\theta})"
        r" = -\,\dfrac{1}{2\,\sigma_{\mathrm{obs}}^2}"
        r"\sum_{k,s}"
        r"\bigl(y_{k,s}"
        r" - \hat{y}_{k,s}(\boldsymbol{\theta})\bigr)^{\!2}"
        r" \;+\; \mathrm{const.}$"
    ),
    prior=r"$p(\boldsymbol{\theta}) = \prod_{j=1}^{20}\mathrm{Uniform}(\theta_j;\; l_j,\, u_j)$",
    rom_approx=(
        r"$\mathbf{g}(t) \approx V\,\mathbf{z}(t),"
        r"\qquad V \in \mathbb{R}^{12\times r},"
        r"\quad \dim\mathbf{z} = r \ll 12$"
    ),
    rom_map=(r"$\mathbf{z}_{n+1}" r" = F_n(\mathbf{z}_n;\,\boldsymbol{\theta})$"),
    linearization_update=(
        r"$\boldsymbol{\theta}_0^{(m+1)}"
        r" \;\leftarrow\; \hat{\boldsymbol{\theta}}_{\mathrm{MAP}}^{(m)}"
        r" \;=\; \arg\max_j\; p_{\beta_m}"
        r"(\boldsymbol{\theta}_j \mid \mathbf{y})$"
    ),
    tempered=(
        r"$p_{\beta}(\boldsymbol{\theta}\mid\mathbf{y})"
        r"\;\propto\;"
        r"p(\mathbf{y}\mid\boldsymbol{\theta})^{\,\beta}\;"
        r"p(\boldsymbol{\theta}),"
        r"\qquad 0 = \beta_0 < \beta_1 < \cdots < \beta_M = 1$"
    ),
    ess=(
        r"$\mathrm{ESS} = \dfrac{1}{\sum_{j=1}^{N}\bar{w}_j^{\,2}}"
        r"\;,\qquad \bar{w}_j = \dfrac{w_j}{\sum_l w_l}$"
    ),
    weight_update=(
        r"$w_j^{(m)} = p(\mathbf{y}\mid\boldsymbol{\theta}_j)^{\," r"\beta_{m+1} - \beta_m}$"
    ),
    proposal_cov=(
        r"$q(\boldsymbol{\theta}^*\mid\boldsymbol{\theta}_j)"
        r" = \mathcal{N}\!\left(\boldsymbol{\theta}_j,\;"
        r"\dfrac{2.38^2}{d}\,\boldsymbol{\Sigma}^{(m)}\right)$"
    ),
    mh_accept=(
        r"$\alpha = \min\!\left\{1,\;"
        r"\dfrac{p(\mathbf{y}\mid\boldsymbol{\theta}^*)^{\,\beta}"
        r"\; p(\boldsymbol{\theta}^*)}"
        r"{p(\mathbf{y}\mid\boldsymbol{\theta}_j)^{\,\beta}"
        r"\; p(\boldsymbol{\theta}_j)}\right\}$"
    ),
    evidence=(
        r"$\hat{p}(\mathbf{y}) = \prod_{m=0}^{M-1}"
        r"\left(\frac{1}{N}\sum_{j=1}^{N} w_j^{(m)}\right)$"
    ),
    bayes_factor=(
        r"$B_{10} = \dfrac{p(\mathbf{y}\mid\mathcal{M}_{\mathrm{reduced}})}"
        r"{p(\mathbf{y}\mid\mathcal{M}_{\mathrm{full}})}$"
    ),
    seq_update=(
        r"$\theta_{\mathrm{base},k}"
        r" \;\leftarrow\; \hat{\theta}_k^{\,\mathrm{MAP}}"
        r" \quad \forall\, k \in \mathcal{A}_s$"
    ),
    # ── Biofilm model: Hamilton principle formulation ────────
    hamilton_local=(
        r"$\dfrac{\partial \Psi}{\partial \boldsymbol{\xi}}"
        r" + \dfrac{\partial \Delta_s}{\partial \dot{\boldsymbol{\xi}}}"
        r" + \dfrac{\partial c}{\partial \boldsymbol{\xi}}"
        r" = \mathbf{0}$"
    ),
    xi_def=(
        r"$\boldsymbol{\xi}"
        r" = \begin{pmatrix} \boldsymbol{\phi}"
        r" \\ \boldsymbol{\psi} \end{pmatrix},"
        r"\qquad \bar{\phi}_l = \phi_l\,\psi_l"
        r"\;\;\text{(living bacteria volume)}$"
    ),
    energy_psi=(
        r"$\Psi = -\,\dfrac{1}{2}\,c^{\!*}\,"
        r"\boldsymbol{\bar{\phi}} \cdot A \cdot \boldsymbol{\bar{\phi}}"
        r" \;+\; \dfrac{1}{2}\,\alpha^{\!*}\,"
        r"\boldsymbol{\psi} \cdot B \cdot \boldsymbol{\psi}$"
    ),
    dissipation_delta=(
        r"$\Delta_s"
        r" = \dfrac{1}{2}\,"
        r"\dot{\boldsymbol{\bar{\phi}}} \cdot \boldsymbol{\eta}"
        r" \cdot \dot{\boldsymbol{\bar{\phi}}}"
        r" \;+\; \dfrac{1}{2}\,"
        r"\dot{\boldsymbol{\phi}} \cdot \boldsymbol{\eta}"
        r" \cdot \dot{\boldsymbol{\phi}}$"
    ),
    constraint_gamma=(r"$c = \gamma\!\left(" r"\displaystyle\sum_{l=0}^{n} \phi_l - 1\right) = 0$"),
    # ── Penalty barrier + total energy ─────────────────────────
    penalty_barrier=(
        r"$\Psi_p"
        r" = \displaystyle\sum_{l}\dfrac{K_p}{\xi_l^{\,2}\,(1-\xi_l)^2},"
        r"\qquad \xi_l \in \{\phi_l,\,\psi_l\}$"
    ),
    psi_total=(
        r"$\Psi_{\mathrm{total}}"
        r" = \underbrace{\Psi}_{\text{eq.\,(13)}}"
        r" + \underbrace{\Psi_p}_{\text{penalty}}"
        r"\;\equiv\; F(\mathbf{g})$"
    ),
    # ── State-vector bridge ξ → g ────────────────────────────
    xi_to_g=(
        r"$\boldsymbol{\xi} = (\boldsymbol{\phi},\,\boldsymbol{\psi})^{\!\top}"
        r"\;\longrightarrow\;"
        r"\mathbf{g} = (\phi_1,\dots,\phi_5,\;\phi_0,\;"
        r"\psi_1,\dots,\psi_5,\;\gamma)^{\!\top}$"
    ),
    # ── Strong form (color-coded, 1-slide overview) ──────────
    strong_phi_color=(
        r"$Q_{\phi_i} \;=\; "
        r"\dfrac{K_p(2\!-\!4\phi_i)}{(\phi_i\!-\!1)^3\phi_i^3}"
        r"\;+\;\dfrac{1}{\eta_i}\Bigl["
        r"{\color[rgb]{0.13,0.55,0.13}\gamma}"
        r" + {\color[rgb]{0.80,0.20,0.00}"
        r"(\eta_{\phi,i}\!+\!\eta_i\psi_i^2)\dot{\phi}_i"
        r"\!+\!\eta_i\phi_i\psi_i\dot{\psi}_i}"
        r"\Bigr]"
        r"\;-\;{\color[rgb]{0.00,0.27,0.80}"
        r"\dfrac{c^{\!*}}{\eta_i}\psi_i"
        r"\bigl(A(\boldsymbol{\phi}"
        r"\!\odot\!\boldsymbol{\psi})\bigr)_{\!i}}$"
    ),
    strong_legend=(
        r"$\text{Black: Penalty }\nabla\Psi_p"
        r"\qquad {\color[rgb]{0.80,0.20,0.00}"
        r"\text{Orange: Dissipation}}"
        r"\qquad {\color[rgb]{0.13,0.55,0.13}"
        r"\text{Green: Constraint }\gamma}"
        r"\qquad {\color[rgb]{0.00,0.27,0.80}"
        r"\text{Blue: Interaction}}$"
    ),
    # ── Hamilton principle → implementation ───────────────────
    hamilton_var=(
        r"$\delta \!\int_0^T \!\bigl[\,"
        r"F(\mathbf{g}) + D(\dot{\mathbf{g}},\mathbf{g})"
        r"\,\bigr]\,dt = 0$"
    ),
    implicit_euler_full=(
        r"$\dfrac{\mathbf{g}^{\,n+1} - \mathbf{g}^{\,n}}{\Delta t}"
        r" = -\,M\!\bigl(\mathbf{g}^{\,n+1}\bigr)\;"
        r"\nabla F\!\bigl(\mathbf{g}^{\,n+1}\bigr)"
        r" + S\!\bigl(\mathbf{g}^{\,n+1}\bigr)$"
    ),
    residual_newton=(
        r"$Q\!\bigl(\mathbf{g}^{\,n+1};\,"
        r"\mathbf{g}^{\,n},\,\boldsymbol{\theta}\bigr) = \mathbf{0}"
        r"\quad\xrightarrow{\;\text{Newton}\;}\quad"
        r"\mathbf{g}^{\,n+1}_{k+1} = \mathbf{g}^{\,n+1}_{k}"
        r" - J^{-1}\!Q$"
    ),
)

# ═══════════════════════════════════════════════════════════════
#  Slide content definitions — Japanese
# ═══════════════════════════════════════════════════════════════
SLIDES_JA: list[tuple[str, list]] = [
    # ── 1. Outline ─────────────────────────────────────────────
    (
        "目次  —  Outline",
        [
            (
                "b",
                [
                    "1.   研究の背景と目的",
                    "2.   5菌種と相互作用ネットワーク",
                    "3.   バイオフィルム成長モデルへの適用",
                    "4.   FOM — 状態変数と制約条件",
                    "5.   FOM — 自由エネルギーに基づく力学方程式",
                    "6.   FOM — 強形式の構造（概観）",
                    "7.   FOM — 残差方程式の詳細  (φᵢ の方程式)",
                    "8.   FOM — 残差方程式の詳細  (ψᵢ, γ の方程式)",
                    "9.   パラメータ — 相互作用行列 A と減衰行列 B",
                    "10.  生物学的制約によるパラメータ削減",
                    "11.  ヒル関数によるゲーティング機構",
                    "12.  観測モデルと尤度関数",
                    "13.  ROM — TSM と線形化点更新",
                    "14.  TMCMC — アルゴリズムの詳細",
                    "15.  4段階逐次推定と実験条件",
                    "16.  まとめと今後の展望",
                    "Backup:  Hamilton 原理から実装へ",
                ],
                20,
            ),
        ],
    ),
    # ── 2. Background & objectives ─────────────────────────────
    (
        "研究の背景と目的",
        [
            ("h", "背景"),
            (
                "b",
                [
                    "口腔バイオフィルムは複数の細菌種が共存する複雑な微生物群集である.",
                    "種間の相互作用 (協力・競合) がバイオフィルムの構造と病原性を決定する.",
                    "実験データから種間相互作用を定量的に推定することが臨床応用の鍵となる.",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "目的"),
            (
                "b",
                [
                    "自由エネルギーに基づく 5 種バイオフィルムの力学モデル (FOM) を構築.",
                    "ROM による高速化と TMCMC によるベイズ推定を組み合わせて,",
                    "種間相互作用パラメータ θ を実験データから確率的に推定する.",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "アプローチ (3 段階)"),
            (
                "b",
                [
                    "1.  FOM: 自由エネルギーに基づく非線形散逸力学系 → 高精度だが計算コスト大",
                    "2.  ROM: FOM を低次元近似 → forward 評価を高速化",
                    "3.  TMCMC: ROM で尤度を高速評価 → 事後分布 p(θ|y) をサンプリング",
                ],
                18,
            ),
        ],
    ),
    # ── 3. Species & Interaction Network (NEW) ─────────────────
    (
        "5菌種と相互作用ネットワーク",
        [
            ("h", "対象菌種 (Heine et al., 2025)"),
            (
                "b",
                [
                    "1.  S. oralis  — 初期定着菌 (early colonizer)",
                    "2.  A. naeslundii  — 初期定着菌 (early colonizer)",
                    "3.  Veillonella spp.  — 代謝橋 (metabolic bridge, 乳酸消費)",
                    "4.  F. nucleatum  — 橋渡し菌 (bridge organism)",
                    "5.  P. gingivalis  — 後期定着菌 (late colonizer, 歯周病原菌)",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "相互作用ネットワーク (Figure 4C)"),
            (
                "b",
                [
                    "S.o ↔ A.n:  共凝集 (co-aggregation)",
                    "S.o ↔ Vei:   乳酸の供給・消費 (lactate handover)",
                    "S.o ↔ F.n:   ギ酸・酢酸の共生 (formate/acetate symbiosis)",
                    "Vei ↔ P.g:    pH上昇による支援 (pH rise support)",
                    "F.n ↔ P.g:    共凝集・ペプチド供給 (co-aggregation, peptides)",
                ],
                17,
            ),
            ("g", 0.10),
            (
                "tc",
                "実験的に存在しない相互作用 (5ペア) はゼロに固定 → パラメータ削減",
                RED_ACCENT,
                17,
                True,
            ),
        ],
    ),
    # ── 3. Application to biofilm growth ─────────────────────
    (
        "バイオフィルム成長モデルへの適用",
        [
            ("t", "Klempt et al. の拡張 Hamilton 原理に基づくバイオフィルムモデル:", 18),
            ("g", 0.08),
            ("h", "拡張 Hamilton 原理 (局所・準静的・等温)"),
            ("e", _EQ["hamilton_local"], 24),
            ("g", 0.05),
            ("e", _EQ["xi_def"], 20),
            ("g", 0.10),
            ("h", "エネルギー密度関数 Ψ  (eq.13)"),
            ("e", _EQ["energy_psi"], 22),
            ("g", 0.05),
            (
                "b",
                [
                    "第1項:  栄養素 c* による生菌増加 — 相互作用行列 A が結合強度を記述",
                    "第2項:  抗菌薬 α* による生菌減少 — 減衰行列 B が感受性を記述",
                ],
                16,
            ),
            ("g", 0.08),
            ("h", "散逸関数 Δs  (eq.16)  と  制約関数 c  (eq.18)"),
            ("e", _EQ["dissipation_delta"], 20),
            ("g", 0.05),
            ("e", _EQ["constraint_gamma"], 22),
            ("g", 0.08),
            ("h", "ペナルティ障壁 Ψ_p  (ξ ∈ [0,1] を保証)"),
            ("e", _EQ["penalty_barrier"], 20),
            ("g", 0.05),
            ("e", _EQ["psi_total"], 20),
            ("g", 0.05),
            (
                "t",
                "→  散逸が φ̄ = φψ の時間微分で定義されるため, φ と ψ が深く結合した系となる.",
                16,
            ),
        ],
    ),
    # ── 4. State variables ─────────────────────────────────────
    (
        "FOM — 状態変数と制約条件",
        [
            ("h", "記法の対応:  ξ → g"),
            ("e", _EQ["xi_to_g"], 16),
            ("t", "φ₀ (マトリックス体積分率) と γ (ラグランジュ乗数) を状態変数に追加.", 16),
            ("g", 0.10),
            ("tb", "状態ベクトル"),
            ("g", 0.08),
            ("e", _EQ["state_vec"]),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$\phi_i(t)$", "種 i のバイオフィルム体積分率  (i = 1, …, 5)"),
                    (r"$\phi_0(t)$", "マトリックス / 水相の体積分率"),
                    (r"$\psi_i(t)$", "種 i の内部構造・密度変数"),
                    (r"$\gamma(t)$", "体積保存を満たすラグランジュ乗数"),
                ],
            ),
            ("g", 0.15),
            ("tb", "体積保存の制約"),
            ("g", 0.08),
            ("e", _EQ["vol_constraint"]),
            ("g", 0.08),
            ("t", "→  全ての種とマトリックスの体積分率の和は常に 1 に保たれる.", 18),
        ],
    ),
    # ── 5. Governing equations ─────────────────────────────────
    (
        "FOM — 自由エネルギーに基づく力学方程式",
        [
            ("tb", "Ψ_total ≡ F(g) に基づく散逸力学"),
            ("g", 0.08),
            ("e", _EQ["dissipative"], 28),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$M(\mathbf{g})$", "モビリティ行列 — φ̄ = φψ 結合を含む (η⁻¹ ではない)"),
                    (
                        r"$F(\mathbf{g}) \equiv \Psi_{\mathrm{total}}$",
                        "全エネルギー = 相互作用 Ψ + ペナルティ Ψ_p",
                    ),
                    (r"$S(\mathbf{g})$", "ソース項 — 成長・洗い流し等の非保存力"),
                ],
            ),
            ("g", 0.15),
            ("tb", "時間離散化 — 暗示的オイラー法"),
            ("g", 0.08),
            ("e", _EQ["implicit_euler"], 28),
            ("g", 0.08),
            ("t", "→  各タイムステップで非線形連立方程式をニュートン法で求解.", 18),
            ("t", "→  硬い系 (stiff system) に対し安定な時間積分を実現.", 18),
        ],
    ),
    # ── 5b. Strong form overview (color-coded) ─────────────────
    (
        "FOM — 強形式の構造（概観）",
        [
            ("t", "体積分率 φᵢ に対する強形式を 3 つの物理的項に色分け:", 20),
            ("g", 0.10),
            ("e", _EQ["strong_phi_color"], 20),
            ("g", 0.15),
            ("e", _EQ["strong_legend"], 14),
            ("g", 0.20),
            ("tb", "体積保存の制約"),
            ("g", 0.05),
            ("e", _EQ["vol_constraint"], 24),
            ("g", 0.15),
            (
                "b",
                [
                    "ペナルティ勾配:  障壁関数 Ψ_p の φᵢ 微分 → ξ ∈ (0,1) を保証",
                    "散逸:  粘性項 (η_φ, ηᵢ) による時間変化の減衰",
                    "制約:  ラグランジュ乗数 γ が体積保存を強制",
                    "相互作用:  c* · A · (φ⊙ψ) による種間結合 (栄養素 c* が駆動)",
                ],
                16,
            ),
        ],
    ),
    # ── 6. Residual phi ────────────────────────────────────────
    (
        "残差方程式 — φᵢ の方程式",
        [
            ("t", "体積分率 φᵢ に対する残差方程式:", 20),
            ("g", 0.08),
            ("e", _EQ["res_phi_1"], 24),
            ("g", 0.08),
            ("e", _EQ["res_phi_2"], 24),
            ("g", 0.08),
            ("e", _EQ["res_phi_3"], 24),
            ("g", 0.20),
            ("h", "各項の物理的意味"),
            (
                "b",
                [
                    "Term 1:  ペナルティ障壁 ∂Ψ_p/∂φᵢ — ξ ∈ (0,1) を保証する障壁力",
                    "Term 2:  散逸項 — ラグランジュ乗数 γ と粘性 (η_φ, ηᵢ) による減衰",
                    "Term 3:  種間相互作用 — 栄養素 c* と行列 A による結合 (φ⊙ψ)",
                ],
                17,
            ),
        ],
    ),
    # ── 7. Residual psi, gamma ─────────────────────────────────
    (
        "残差方程式 — ψᵢ, γ の方程式",
        [
            ("tb", "ψᵢ に対する残差"),
            ("g", 0.08),
            ("e", _EQ["res_psi_1"], 22),
            ("g", 0.08),
            ("e", _EQ["res_psi_2"], 22),
            ("g", 0.15),
            (
                "b",
                [
                    "Terms 1–2:  ペナルティ障壁 ∂Ψ_p/∂ψᵢ — ψ ∈ (0,1) を保証",
                    "Term 3:  減衰項 — 抗菌薬 α* と減衰率 bᵢ による生存率の低下",
                    "Terms 4–5:  体積分率との結合 + 栄養素 c* による種間相互作用",
                ],
                17,
            ),
            ("g", 0.15),
            ("tb", "体積保存の制約残差"),
            ("g", 0.08),
            ("e", _EQ["res_gamma"], 24),
        ],
    ),
    # ── 8. Matrices A, B ───────────────────────────────────────
    (
        "パラメータ — 相互作用行列 A と減衰行列 B",
        [
            ("tb", "推定パラメータベクトル  θ ∈ R²⁰ (5ブロック構成)"),
            ("g", 0.08),
            ("e", _EQ["theta_vec"], 16),
            ("g", 0.20),
            ("mp", _EQ["mat_A"], 20, _EQ["mat_B"], 20),
            ("g", 0.25),
            (
                "b",
                [
                    "A は対称行列 (aᵢⱼ = aⱼᵢ):  上三角 15 独立成分",
                    "aᵢⱼ > 0: 種 i と j の協力関係,  aᵢⱼ < 0: 競合関係",
                    "B は対角行列:  各種の減衰率 bᵢ → 5 パラメータ",
                    "合計 20 パラメータを 5 ブロック (M1–M5) に分割",
                ],
                17,
            ),
        ],
    ),
    # ── 9. Biological Constraints (NEW) ────────────────────────
    (
        "生物学的制約によるパラメータ削減",
        [
            ("h", "実験に基づくパラメータ固定 (Nishioka Algorithm)"),
            (
                "t",
                "Heine et al. の実験 (Figure 4C) で相互作用が観測されない 5 ペアをゼロに固定:",
                18,
            ),
            ("g", 0.10),
            ("e", _EQ["locked_set"], 20),
            ("g", 0.10),
            (
                "b",
                [
                    "a₃₄ = 0:  Vei ↔ F.n — 直接の代謝経路なし",
                    "a₂₃ = 0:  A.n ↔ Vei — 直接の代謝リンクなし",
                    "a₂₄ = 0:  A.n ↔ F.n — 直接の相互作用なし",
                    "a₁₅ = 0:  S.o ↔ P.g — 直接の相互作用なし",
                    "a₂₅ = 0:  A.n ↔ P.g — 直接の相互作用なし",
                ],
                16,
            ),
            ("g", 0.10),
            ("h", "有効パラメータ次元"),
            ("e", _EQ["n_free"], 24),
            ("g", 0.08),
            ("t", "→  サンプリング効率向上 + 生物学的に不合理な推定値を排除.", 18),
        ],
    ),
    # ── 10. Hill gating ────────────────────────────────────────
    (
        "ヒル関数によるゲーティング機構",
        [
            ("t", "F. nucleatum (種 4) が P. gingivalis (種 5) の橋渡し菌として機能する.", 20),
            ("t", "橋渡し菌が十分に存在しないと, P. gingivalis の相互作用は抑制される.", 18),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$x_{\mathrm{FN}} = \phi_4\,\psi_4$", "F. nucleatum の有効バイオマス"),
                ],
            ),
            ("g", 0.10),
            ("tb", "ヒル関数"),
            ("g", 0.05),
            ("e", _EQ["hill"], 28),
            ("g", 0.08),
            (
                "s",
                [
                    (r"$K_{\mathrm{hill}}$", "半飽和定数 — ゲートが 50% ON になる濃度"),
                    (r"$n_{\mathrm{hill}}$", "Hill 係数 — 応答の急峻度を制御"),
                ],
                18,
            ),
            ("g", 0.10),
            ("tb", "有効相互作用"),
            ("g", 0.05),
            ("e", _EQ["a_eff"], 24),
        ],
    ),
    # ── 11. Observation model ──────────────────────────────────
    (
        "観測モデルと尤度関数",
        [
            ("tb", "観測モデル"),
            ("g", 0.08),
            ("e", _EQ["obs_model"], 24),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$y_{k,s}$", "時刻 t_k, 種 s での観測値"),
                    (r"$\hat{y}_{k,s}(\boldsymbol{\theta})$", "ROM を用いたモデル予測値"),
                    (r"$\sigma_{\mathrm{obs}}$", "観測ノイズの標準偏差"),
                ],
            ),
            ("g", 0.10),
            ("tb", "対数尤度"),
            ("g", 0.08),
            ("e", _EQ["loglik"], 24),
            ("g", 0.10),
            ("tb", "事前分布"),
            ("g", 0.08),
            ("e", _EQ["prior"], 22),
        ],
    ),
    # ── 12. ROM / TSM + Linearization ──────────────────────────
    (
        "ROM — TSM と線形化点更新",
        [
            ("tb", "次元削減"),
            ("g", 0.08),
            ("e", _EQ["rom_approx"], 22),
            ("g", 0.08),
            ("t", "V はスナップショットの POD (Proper Orthogonal Decomposition) から構成.", 18),
            ("g", 0.12),
            ("tb", "FOM の線形化 → 時間発展マップ"),
            ("g", 0.08),
            ("e", _EQ["rom_map"], 26),
            ("g", 0.12),
            ("h", "TMCMC 中の線形化点更新"),
            ("t", "TMCMC の各ステージで MAP 推定値に基づき線形化点 θ₀ を更新:", 17),
            ("g", 0.05),
            ("e", _EQ["linearization_update"], 20),
            ("g", 0.08),
            (
                "b",
                [
                    "事後分布が集中するにつれ TSM の近似精度が段階的に向上",
                    "ROM 誤差が閾値未満なら更新をスキップ → FOM 評価コストを削減",
                ],
                17,
            ),
        ],
    ),
    # ── 13. TMCMC Details ──────────────────────────────────────
    (
        "TMCMC — アルゴリズムの詳細",
        [
            ("tb", "テンパリングされた事後分布"),
            ("g", 0.05),
            ("e", _EQ["tempered"], 20),
            ("g", 0.10),
            ("tc", "1)  重み計算", ACCENT, 17, True),
            ("e", _EQ["weight_update"], 18),
            ("g", 0.05),
            ("tc", "2)  適応的 β スケジュール", ACCENT, 17, True),
            ("t", "     ESS が目標値 (N/2) を維持するよう β_{m+1} を二分法で決定:", 16),
            ("e", _EQ["ess"], 18),
            ("g", 0.05),
            ("tc", "3)  リサンプリング + MCMC 変異", ACCENT, 17, True),
            ("t", "     提案分布 (Gelman 最適スケーリング 2.38²/d):", 16),
            ("e", _EQ["proposal_cov"], 18),
            ("g", 0.05),
            ("tc", "4)  MH 受容確率", ACCENT, 17, True),
            ("e", _EQ["mh_accept"], 17),
        ],
    ),
    # ── 14. Sequential Estimation + Conditions (NEW) ───────────
    (
        "4段階逐次推定と実験条件",
        [
            ("h", "4段階逐次推定"),
            (
                "b",
                [
                    "Stage 1 (M1):  S.o + A.n  → 5 params (a₁₁, a₁₂, a₂₂, b₁, b₂)",
                    "Stage 2 (M2):  Vei + F.n  → 5 params (a₃₃, a₃₄, a₄₄, b₃, b₄)",
                    "Stage 3 (M3+M4):  Cross + P.g self → 6 params (a₁₃, a₁₄, a₂₃, a₂₄, a₅₅, b₅)",
                    "Stage 4 (M5):  P.g cross → 4 params (a₁₅, a₂₅, a₃₅, a₄₅)",
                ],
                16,
            ),
            ("g", 0.05),
            ("t", "各ステージ完了後, MAP 推定値を θ_base に固定して次ステージへ:", 17),
            ("e", _EQ["seq_update"], 20),
            ("g", 0.10),
            ("h", "4 実験条件 (Heine et al.)"),
            (
                "b",
                [
                    "Commensal Static:  ロック 9 / 推定 11 — 病原菌ゼロの基本条件",
                    "Dysbiotic Static:   ロック 5 / 推定 15 — 病原菌あり・静的培養",
                    "Commensal HOBIC:  ロック 8 / 推定 12 — フロー培養・S.o 増殖推定",
                    "Dysbiotic HOBIC:   ロック 0 / 推定 20 — 全パラメータ解放 (Surge 再現)",
                ],
                16,
            ),
            ("g", 0.05),
            ("t", "→  条件ごとにロック集合を変え, 生物学的に妥当な推定空間を構成.", 17),
        ],
    ),
    # ── 15. Summary ────────────────────────────────────────────
    (
        "まとめと今後の展望",
        [
            ("h", "まとめ"),
            (
                "b",
                [
                    "FOM:  自由エネルギーに基づく 5 種バイオフィルムの非線形散逸力学系を構築.",
                    "パラメータ削減: 生物学的制約で 20 → 15 自由パラメータに削減.",
                    "ROM (TSM):  FOM を低次元空間で近似 + 線形化点を適応的に更新.",
                    "TMCMC:  適応的 β スケジュール・最適提案共分散・K-step 変異で効率的推定.",
                    "4段階逐次推定:  生物学的グループに基づく分割で安定したベイズ推定.",
                ],
                17,
            ),
            ("g", 0.15),
            ("h", "今後の展望"),
            (
                "b",
                [
                    "実験データへの適用と推定精度の検証.",
                    "モデルエビデンス: TMCMCが自然に提供するベイズファクターでモデルを比較.",
                    "不確実性定量化:  事後分布から予測の信頼区間を定量的に評価.",
                    "多条件推定:  4実験条件 (commensal/dysbiotic × static/HOBIC) に拡張.",
                ],
                17,
            ),
            ("g", 0.10),
            ("h", "モデルエビデンス (TMCMC の副産物)"),
            ("e", _EQ["evidence"], 20),
        ],
    ),
    # ── Backup: Hamilton → implementation ──────────────────────
    (
        "バックアップ：Hamilton 原理から実装へ",
        [
            ("h", "Step 1:  変分原理 (Hamilton)"),
            ("e", _EQ["hamilton_var"], 24),
            ("g", 0.10),
            ("h", "Step 2:  強形式 (Euler–Lagrange)"),
            ("e", _EQ["dissipative"], 22),
            ("g", 0.10),
            ("h", "Step 3:  時間離散化 (暗示的オイラー法)"),
            ("e", _EQ["implicit_euler_full"], 20),
            ("g", 0.10),
            ("h", "Step 4:  残差方程式 + Newton 反復"),
            ("e", _EQ["residual_newton"], 18),
            ("g", 0.15),
            ("t", "→  紙の eq.(16)–(18) がこの 4 段階に対応する.", 17),
            ("t", "→  Step 4 の Q が前スライドで示した色分け方程式に相当.", 17),
        ],
    ),
]

# ═══════════════════════════════════════════════════════════════
#  Speaker notes — Japanese
# ═══════════════════════════════════════════════════════════════
NOTES_JA: list[str] = [
    # 1. Outline
    (
        "このスライドでは，全体の流れを俯瞰してお見せします。"
        "前半で自由エネルギーに基づくFOMの構造を説明し，"
        "その後，生物学的制約によるパラメータ削減，"
        "ROMによる次元削減とTMCMCによるベイズ推定，"
        "最後に4段階逐次推定と実験条件の話へと進みます。"
    ),
    # 2. Background
    (
        "ここでは，口腔バイオフィルムを対象とした研究背景と，"
        "本研究の具体的な目的を整理しています。"
        "臨床応用上重要なのは，単に予測するだけでなく，"
        "種間相互作用をパラメータとして定量化する点です。"
        "FOM→ROM→TMCMCという三段階のアプローチを提示します。"
    ),
    # 3. Species & Network (NEW)
    (
        "このスライドでは，5種の細菌名と生物学的役割を具体的に示します。"
        "S. oralis と A. naeslundii が初期定着菌，Veillonella が代謝橋，"
        "F. nucleatum が橋渡し菌，P. gingivalis が病原菌です。"
        "相互作用ネットワークは Heine et al. (2025) の Figure 4C に基づいており，"
        "5つの活性ペアと5つの不在ペアが実験的に確認されています。"
        "この不在ペアの情報がパラメータ削減の根拠となります。"
    ),
    # 3. Application to biofilm growth
    (
        "Klempt et al.のバイオフィルムモデルの理論的基盤を示します。"
        "拡張Hamilton原理から出発し，局所・準静的・等温の特殊ケースとして"
        "エネルギー密度Ψ，散逸関数Δs，制約関数cの3つの構成関数を定義しています。"
        "栄養素c*が生菌の増加を促進し，抗菌薬α*が減少をもたらすという物理が"
        "行列AとBに集約されています。"
        "散逸関数がφ̄=φψの時間微分で定義される点が，"
        "このモデルの複雑な結合構造の起源です。"
    ),
    # 4. State variables
    (
        "FOMの状態ベクトルg(t)の構成を示します。"
        "5種の体積分率φ_iに加えて，マトリックスφ_0，内部構造変数ψ_i，"
        "体積保存を保証するラグランジュ乗数γをまとめて扱っています。"
        "体積分率の総和が常に1になる制約がモデルの基本構造を規定しています。"
    ),
    # 5. Governing equations
    (
        "全エネルギーΨ_total≡F(g)に基づく散逸力学系としてFOMを定式化している点を説明します。"
        "モビリティMはφ̄=φψ結合を含むため単純なη⁻¹ではない点に注意してください。"
        "S=0のとき全エネルギーFは時間とともに単調減少し，熱力学的に整合的です。"
        "数値的には暗示オイラー法とニュートン法を組み合わせ，"
        "硬いダイナミクスに対しても安定な時間積分を実現しています。"
    ),
    # 5b. Strong form overview
    (
        "このスライドでは，φᵢ の強形式を1枚にまとめ，"
        "各項の物理的役割を色で直感的に示しています。"
        "黒がペナルティ障壁∇Ψ_pの勾配，"
        "オレンジが粘性による散逸，"
        "緑がラグランジュ乗数γによる体積保存制約，"
        "青が行列Aを介した種間相互作用です。"
        "次のスライドで各項の数学的詳細を展開します。"
    ),
    # 6. Residual phi
    (
        "体積分率φ_iに対する残差Q_{φ_i}の具体的な形を示します。"
        "第1項がペナルティ障壁Ψ_pの勾配，第2項が粘性やラグランジュ乗数による散逸，"
        "第3項が行列Aを介した種間相互作用を表しています。"
    ),
    # 7. Residual psi, gamma
    (
        "内部変数ψ_iおよび制約用のγに対する残差方程式をまとめています。"
        "ψ_iについては，ポテンシャル勾配に加えて，減衰率b_iによる緩和項と，"
        "φ_iとの結合項が現れます。"
        "Q_γの式で体積保存条件が方程式系に組み込まれています。"
    ),
    # 8. Matrices A, B
    (
        "推定対象のパラメータベクトルθと，行列A, Bの構造を示します。"
        "θは5ブロック (M1–M5) に分割されており，"
        "これが後で説明する4段階逐次推定のステージ分割に対応します。"
        "Aは対称行列で15個の独立成分，Bは対角行列で5成分，合計20パラメータです。"
    ),
    # 9. Biological Constraints (NEW)
    (
        "これがNishioka Algorithmの核心部分です。"
        "実験的に相互作用が確認されていない5ペアのパラメータを"
        "ゼロに固定することで，推定空間を20次元から15次元に削減します。"
        "これにより，サンプリング効率が向上し，生物学的に不合理な"
        "パラメータ推定値を事前に排除できます。"
        "条件によってはロック数が変わる点は後のスライドで説明します。"
    ),
    # 10. Hill gating
    (
        "F. nucleatumがP. gingivalisの橋渡し菌として働くという"
        "生物学的知見を，ヒル関数h(x_FN)としてモデルに埋め込んでいます。"
        "x_FNはF. nucleatumの有効バイオマスであり，"
        "その大きさに応じてP. gingivalisに関する行列Aの第5行がスイッチされます。"
    ),
    # 11. Observation model
    (
        "観測モデルと対数尤度の定義をまとめています。"
        "観測値y_{k,s}はROMによる予測値にガウスノイズを加えたものとみなし，"
        "二乗和に基づく標準的なガウス尤度を採用しています。"
        "事前分布は各成分独立な一様分布です。"
    ),
    # 12. ROM / TSM + Linearization
    (
        "FOMをPOD基底Vで次元削減し，TSMとして時間構造に沿った更新を実現しています。"
        "重要な改良点として，TMCMC中に線形化点θ₀をMAP推定値に基づいて"
        "適応的に更新することで，事後分布が集中するにつれてROMの近似精度が向上します。"
        "ROM誤差が小さい場合は更新をスキップしてFOM評価コストを削減します。"
    ),
    # 13. TMCMC Details
    (
        "TMCMCのアルゴリズム詳細を4ステップに分けて説明します。"
        "βのスケジュールはESSに基づく二分法で適応的に決定されます。"
        "提案共分散にはGelman et al. (1996) の最適スケーリング 2.38²/d を使用し，"
        "前のステージの受容率に基づいて適応的にスケールを調整します。"
        "K-step MH変異により，リサンプリング後の粒子相関を低減しています。"
    ),
    # 14. Sequential + Conditions (NEW)
    (
        "4段階逐次推定は，本研究の鍵となるアルゴリズム設計です。"
        "M1 (初期定着菌ペア) → M2 (橋渡し菌ペア) → M3+M4 (交差項+P.g自己) → M5 (P.g交差) "
        "の順に推定します。各ステージの完了後，MAP推定値をθ_baseに固定して次へ進みます。"
        "4実験条件ではロック数が異なります。"
        "特にDysbiotic HOBICではロックをすべて解放し，20パラメータ全体を推定して"
        "Surge現象の再現に必要な最小限の相互作用を発見します。"
    ),
    # 15. Summary
    (
        "最後に，研究のポイントをまとめます。"
        "FOMの構造，生物学的制約によるパラメータ削減，ROM+線形化点更新，"
        "適応的TMCMCアルゴリズム，4段階逐次推定の5つが主要な貢献です。"
        "今後はモデルエビデンスによるモデル比較や，"
        "4実験条件への系統的な適用を計画しています。"
        "モデルエビデンスの式は，TMCMCの各ステージの平均重みの積として"
        "自然に得られる副産物であり，追加計算なしにベイズファクターが計算できます。"
    ),
    # Backup: Hamilton → implementation
    (
        "このバックアップスライドでは，紙のeq.(16)–(18)に対応する"
        "4段階の数学的構造を示しています。"
        "Step 1: Hamilton原理の変分形式から出発し，"
        "Step 2: Euler–Lagrangeによる強形式を導出，"
        "Step 3: 暗示的オイラー法で時間離散化，"
        "Step 4: 残差Q=0としてNewton法で求解します。"
        "前のスライドの色分け方程式はStep 4のQの具体的な形に対応します。"
    ),
]

# ═══════════════════════════════════════════════════════════════
#  Speaker notes — English
# ═══════════════════════════════════════════════════════════════
NOTES_EN: list[str] = [
    # 1. Outline
    (
        "This slide provides an overview of the talk. "
        "We start with the FOM structure based on free energy, "
        "then introduce biological constraints and parameter reduction, "
        "ROM for speed-up, TMCMC for Bayesian estimation, "
        "and finally the 4-stage sequential estimation with experimental conditions."
    ),
    # 2. Background
    (
        "Here we outline the research background and objectives. "
        "Oral biofilms involve multiple species whose interactions determine pathogenicity. "
        "The key challenge is to quantitatively estimate these interactions from experimental data. "
        "Our approach combines FOM, ROM, and TMCMC in three stages."
    ),
    # 3. Species & Network (NEW)
    (
        "This slide introduces the five bacterial species and their biological roles. "
        "S. oralis and A. naeslundii are early colonizers, Veillonella is a metabolic bridge, "
        "F. nucleatum is the bridge organism, and P. gingivalis is the late colonizer (pathogen). "
        "The interaction network is based on Heine et al. (2025) Figure 4C. "
        "Five active interaction pairs and five absent pairs have been experimentally confirmed. "
        "The absent pairs form the basis for our parameter reduction."
    ),
    # 3. Application to biofilm growth
    (
        "This slide presents the theoretical foundation of the biofilm model by Klempt et al. "
        "Starting from the extended Hamilton principle for a local, quasi-static, isothermal system, "
        "we define three constitutive functions: energy density Psi, dissipation Delta_s, and constraint c. "
        "Nutrients c-star promote growth through the interaction matrix A, "
        "while antibiotics alpha-star cause decay through matrix B. "
        "The key complexity arises from modelling dissipation via the rate of living bacteria volume "
        "phi-bar equals phi times psi, which deeply couples the phi and psi equations."
    ),
    # 4. State variables
    (
        "We present the FOM state vector g(t) with 12 components: "
        "five species volume fractions, matrix volume fraction, "
        "five internal structure variables, and a Lagrange multiplier for volume conservation."
    ),
    # 5. Governing equations
    (
        "The FOM is formulated as a dissipative system driven by free energy F(g). "
        "The mobility matrix M encodes the phi-bar coupling, F is the total energy Psi_total, "
        "and S captures growth and washout. When S=0, F decreases monotonically, "
        "ensuring thermodynamic consistency. "
        "We use implicit Euler with Newton's method for stable time integration."
    ),
    # 5b. Strong form overview
    (
        "This slide presents the strong form for phi_i in a single colour-coded equation. "
        "Black is the penalty barrier gradient, "
        "orange highlights the viscous dissipation terms, "
        "green marks the Lagrange multiplier gamma for volume conservation, "
        "and blue shows the inter-species interaction via the matrix A. "
        "The next slides expand each term in mathematical detail."
    ),
    # 6. Residual phi
    (
        "This slide shows the residual equation for the volume fraction phi_i. "
        "Term 1 is the penalty barrier gradient (ensures xi stays in (0,1)), "
        "Term 2 is dissipation (Lagrange multiplier and viscous damping), "
        "Term 3 is the inter-species interaction via matrix A with Hadamard product."
    ),
    # 7. Residual psi, gamma
    (
        "Here we show the residual equations for the internal variable psi_i and "
        "the volume conservation constraint gamma. "
        "The decay rate b_i drives relaxation of the internal variable."
    ),
    # 8. Matrices A, B
    (
        "The parameter vector theta is organized into 5 blocks (M1 through M5). "
        "This block structure directly corresponds to the 4-stage sequential estimation. "
        "A is symmetric with 15 independent entries, B is diagonal with 5 entries, totalling 20 parameters."
    ),
    # 9. Biological Constraints (NEW)
    (
        "This is the core of the Nishioka Algorithm. "
        "We lock 5 interaction parameters to zero based on absent biological interactions "
        "confirmed by Heine et al. (2025). "
        "This reduces the effective dimension from 20 to 15, "
        "improving sampling efficiency and excluding biologically implausible estimates. "
        "The lock set varies by experimental condition, as explained later."
    ),
    # 10. Hill gating
    (
        "The Hill function encodes the biological knowledge that F. nucleatum "
        "acts as a bridge organism for P. gingivalis. "
        "The effective biomass x_FN gates the 5th row of the interaction matrix A."
    ),
    # 11. Observation model
    (
        "We define the observation model with additive Gaussian noise "
        "and the corresponding log-likelihood. "
        "The prior is independent uniform on each parameter component."
    ),
    # 12. ROM / TSM + Linearization
    (
        "The ROM uses POD-based dimension reduction and time-structured mapping. "
        "A key improvement is the adaptive linearization point update during TMCMC: "
        "after each stage, we update theta_0 to the current MAP estimate, "
        "progressively improving ROM accuracy as the posterior concentrates. "
        "ROM error checks determine whether FOM evaluation is needed."
    ),
    # 13. TMCMC Details
    (
        "We present the full TMCMC algorithm in four steps. "
        "The beta schedule is adaptively determined via bisection to maintain ESS at N/2. "
        "The proposal covariance uses Gelman et al. optimal scaling 2.38^2/d "
        "with tempered scaling and adaptive adjustment based on acceptance rate. "
        "K-step MH mutation reduces particle correlation after resampling."
    ),
    # 14. Sequential + Conditions (NEW)
    (
        "The 4-stage sequential estimation is a key algorithmic design. "
        "Stage 1 estimates early colonizer parameters, Stage 2 bridge organisms, "
        "Stage 3 cross-terms plus P. gingivalis self-interaction, Stage 4 P. gingivalis cross-terms. "
        "After each stage, MAP estimates are fixed in theta_base. "
        "The four experimental conditions vary the lock set: "
        "Commensal Static locks 9 parameters, Dysbiotic HOBIC unlocks all 20 "
        "to discover the minimal interaction set driving the pathogen surge."
    ),
    # 15. Summary
    (
        "We summarize the five main contributions: FOM structure, parameter reduction, "
        "ROM with adaptive linearization, TMCMC with adaptive beta and optimal proposal, "
        "and 4-stage sequential estimation. "
        "The model evidence formula is a natural byproduct of TMCMC, "
        "enabling Bayesian model comparison via Bayes factors at no extra cost."
    ),
    # Backup: Hamilton → implementation
    (
        "This backup slide shows the four-step mathematical correspondence "
        "from the Hamilton variational principle to the Newton-based residual solver. "
        "Step 1: Hamilton's principle gives the variational form. "
        "Step 2: Euler-Lagrange equations yield the strong form. "
        "Step 3: Implicit Euler discretisation in time. "
        "Step 4: The residual Q=0 is solved by Newton iteration. "
        "The colour-coded equation from the earlier slide is the explicit form of Q."
    ),
]

# ═══════════════════════════════════════════════════════════════
#  Slide content definitions — English
# ═══════════════════════════════════════════════════════════════
SLIDES_EN: list[tuple[str, list]] = [
    # ── 1. Outline ─────────────────────────────────────────────
    (
        "Outline",
        [
            (
                "b",
                [
                    "1.   Background and Objectives",
                    "2.   5 Species and Interaction Network",
                    "3.   Application to Biofilm Growth",
                    "4.   FOM — State Variables and Constraints",
                    "5.   FOM — Free-Energy-Based Governing Equations",
                    "6.   FOM — Strong Form Structure (Overview)",
                    "7.   FOM — Residual Equations (φᵢ)",
                    "8.   FOM — Residual Equations (ψᵢ, γ)",
                    "9.   Parameters — Interaction Matrix A and Decay Matrix B",
                    "10.  Biological Constraints and Parameter Reduction",
                    "11.  Hill-Function Gating Mechanism",
                    "12.  Observation Model and Likelihood",
                    "13.  ROM — TSM with Linearization Update",
                    "14.  TMCMC — Algorithm Details",
                    "15.  4-Stage Sequential Estimation and Experimental Conditions",
                    "16.  Summary and Outlook",
                    "Backup:  From Hamilton Principle to Implementation",
                ],
                20,
            ),
        ],
    ),
    # ── 2. Background & objectives ─────────────────────────────
    (
        "Background and Objectives",
        [
            ("h", "Background"),
            (
                "b",
                [
                    "Oral biofilms are complex microbial communities with multiple co-existing species.",
                    "Inter-species interactions (cooperation / competition) govern biofilm structure and pathogenicity.",
                    "Quantitative estimation of these interactions from experimental data is key to clinical applications.",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "Objectives"),
            (
                "b",
                [
                    "Construct a free-energy-based full-order model (FOM) for 5-species biofilm dynamics.",
                    "Combine ROM speed-up with TMCMC Bayesian estimation to infer the",
                    "inter-species interaction parameter θ from experimental data in a probabilistic framework.",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "Approach (3 stages)"),
            (
                "b",
                [
                    "1.  FOM: nonlinear dissipative dynamics based on free energy → accurate but expensive",
                    "2.  ROM: low-dimensional approximation of FOM → fast forward evaluation",
                    "3.  TMCMC: fast likelihood via ROM → sample the posterior p(θ|y)",
                ],
                18,
            ),
        ],
    ),
    # ── 3. Species & Interaction Network (NEW) ─────────────────
    (
        "5 Species and Interaction Network",
        [
            ("h", "Species (Heine et al., 2025)"),
            (
                "b",
                [
                    "1.  S. oralis  — early colonizer",
                    "2.  A. naeslundii  — early colonizer",
                    "3.  Veillonella spp.  — metabolic bridge (lactate consumer)",
                    "4.  F. nucleatum  — bridge organism",
                    "5.  P. gingivalis  — late colonizer (periodontal pathogen)",
                ],
                18,
            ),
            ("g", 0.15),
            ("h", "Interaction network (Figure 4C)"),
            (
                "b",
                [
                    "S.o ↔ A.n:  co-aggregation",
                    "S.o ↔ Vei:   lactate production / consumption (lactate handover)",
                    "S.o ↔ F.n:   formate / acetate symbiosis",
                    "Vei ↔ P.g:    pH rise support",
                    "F.n ↔ P.g:    co-aggregation and peptide supply",
                ],
                17,
            ),
            ("g", 0.10),
            (
                "tc",
                "5 species pairs with no observed interaction → locked to zero (parameter reduction)",
                RED_ACCENT,
                17,
                True,
            ),
        ],
    ),
    # ── 3. Application to biofilm growth ─────────────────────
    (
        "Application to Biofilm Growth",
        [
            ("t", "Biofilm model based on the extended Hamilton principle (Klempt et al.):", 18),
            ("g", 0.08),
            ("h", "Extended Hamilton principle (local, quasi-static, isothermal)"),
            ("e", _EQ["hamilton_local"], 24),
            ("g", 0.05),
            ("e", _EQ["xi_def"], 20),
            ("g", 0.10),
            ("h", "Energy density function Ψ  (eq.13)"),
            ("e", _EQ["energy_psi"], 22),
            ("g", 0.05),
            (
                "b",
                [
                    "Term 1:  nutrients c* promote living bacteria growth — matrix A describes coupling",
                    "Term 2:  antibiotics α* reduce viability — matrix B describes susceptibility",
                ],
                16,
            ),
            ("g", 0.08),
            ("h", "Dissipation Δs  (eq.16)  and  constraint c  (eq.18)"),
            ("e", _EQ["dissipation_delta"], 20),
            ("g", 0.05),
            ("e", _EQ["constraint_gamma"], 22),
            ("g", 0.08),
            ("h", "Penalty barrier Ψ_p  (ensures ξ ∈ [0,1])"),
            ("e", _EQ["penalty_barrier"], 20),
            ("g", 0.05),
            ("e", _EQ["psi_total"], 20),
            ("g", 0.05),
            (
                "t",
                "→  Dissipation defined via φ̄ = φψ rate leads to deeply coupled φ-ψ system.",
                16,
            ),
        ],
    ),
    # ── 4. State variables ─────────────────────────────────────
    (
        "FOM — State Variables and Constraints",
        [
            ("h", "Notation bridge:  ξ → g"),
            ("e", _EQ["xi_to_g"], 16),
            ("t", "Augmented with φ₀ (matrix volume fraction) and γ (Lagrange multiplier).", 16),
            ("g", 0.10),
            ("tb", "State vector"),
            ("g", 0.08),
            ("e", _EQ["state_vec"]),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$\phi_i(t)$", "Volume fraction of species i  (i = 1, …, 5)"),
                    (r"$\phi_0(t)$", "Volume fraction of the matrix / aqueous phase"),
                    (r"$\psi_i(t)$", "Internal structure / density variable of species i"),
                    (r"$\gamma(t)$", "Lagrange multiplier enforcing volume conservation"),
                ],
            ),
            ("g", 0.15),
            ("tb", "Volume conservation constraint"),
            ("g", 0.08),
            ("e", _EQ["vol_constraint"]),
            ("g", 0.08),
            ("t", "→  The sum of all volume fractions is always unity.", 18),
        ],
    ),
    # ── 5. Governing equations ─────────────────────────────────
    (
        "FOM — Free-Energy-Based Governing Equations",
        [
            ("tb", "Dissipative dynamics driven by Ψ_total ≡ F(g)"),
            ("g", 0.08),
            ("e", _EQ["dissipative"], 28),
            ("g", 0.15),
            (
                "s",
                [
                    (
                        r"$M(\mathbf{g})$",
                        "Mobility matrix — encodes φ̄ = φψ coupling (not simply η⁻¹)",
                    ),
                    (
                        r"$F(\mathbf{g}) \equiv \Psi_{\mathrm{total}}$",
                        "Total energy = interaction Ψ + penalty Ψ_p",
                    ),
                    (
                        r"$S(\mathbf{g})$",
                        "Source term — growth, washout, and other non-conservative forces",
                    ),
                ],
            ),
            ("g", 0.15),
            ("tb", "Time discretisation — implicit Euler method"),
            ("g", 0.08),
            ("e", _EQ["implicit_euler"], 28),
            ("g", 0.08),
            ("t", "→  Nonlinear system solved by Newton's method at each time step.", 18),
            ("t", "→  Ensures stable integration for the stiff biofilm dynamics.", 18),
        ],
    ),
    # ── 5b. Strong form overview (color-coded) ─────────────────
    (
        "FOM — Strong Form Structure (Overview)",
        [
            ("t", "Strong form for volume fraction φᵢ — colour-coded by physical role:", 20),
            ("g", 0.10),
            ("e", _EQ["strong_phi_color"], 20),
            ("g", 0.15),
            ("e", _EQ["strong_legend"], 14),
            ("g", 0.20),
            ("tb", "Volume conservation constraint"),
            ("g", 0.05),
            ("e", _EQ["vol_constraint"], 24),
            ("g", 0.15),
            (
                "b",
                [
                    "Penalty gradient:  ∂Ψ_p/∂φᵢ — barrier ensuring ξ ∈ (0,1)",
                    "Dissipation:  viscous damping via (η_φ, ηᵢ)",
                    "Constraint:  Lagrange multiplier γ enforces volume conservation",
                    "Interaction:  nutrients c* × matrix A × (φ⊙ψ) inter-species coupling",
                ],
                16,
            ),
        ],
    ),
    # ── 6. Residual phi ────────────────────────────────────────
    (
        "Residual Equations — φᵢ",
        [
            ("t", "Residual equation for the volume fraction φᵢ :", 20),
            ("g", 0.08),
            ("e", _EQ["res_phi_1"], 24),
            ("g", 0.08),
            ("e", _EQ["res_phi_2"], 24),
            ("g", 0.08),
            ("e", _EQ["res_phi_3"], 24),
            ("g", 0.20),
            ("h", "Physical interpretation of each term"),
            (
                "b",
                [
                    "Term 1:  Penalty barrier ∂Ψ_p/∂φᵢ — ensures ξ ∈ (0,1)",
                    "Term 2:  Dissipation — Lagrange multiplier γ and viscosity (η_φ, ηᵢ)",
                    "Term 3:  Inter-species interaction — nutrients c* and matrix A (φ⊙ψ)",
                ],
                17,
            ),
        ],
    ),
    # ── 7. Residual psi, gamma ─────────────────────────────────
    (
        "Residual Equations — ψᵢ, γ",
        [
            ("tb", "Residual for ψᵢ"),
            ("g", 0.08),
            ("e", _EQ["res_psi_1"], 22),
            ("g", 0.08),
            ("e", _EQ["res_psi_2"], 22),
            ("g", 0.15),
            (
                "b",
                [
                    "Terms 1–2:  Penalty barrier ∂Ψ_p/∂ψᵢ — ensures ψ ∈ (0,1)",
                    "Term 3:  Decay — antibiotics α* and rate bᵢ reduce viability",
                    "Terms 4–5:  Coupling with φᵢ + nutrients c* inter-species interaction via A",
                ],
                17,
            ),
            ("g", 0.15),
            ("tb", "Volume conservation constraint residual"),
            ("g", 0.08),
            ("e", _EQ["res_gamma"], 24),
        ],
    ),
    # ── 8. Matrices A, B ───────────────────────────────────────
    (
        "Parameters — Interaction Matrix A and Decay Matrix B",
        [
            ("tb", "Parameter vector  θ ∈ R²⁰ (5-block structure)"),
            ("g", 0.08),
            ("e", _EQ["theta_vec"], 16),
            ("g", 0.20),
            ("mp", _EQ["mat_A"], 20, _EQ["mat_B"], 20),
            ("g", 0.25),
            (
                "b",
                [
                    "A is symmetric (aᵢⱼ = aⱼᵢ):  15 independent upper-triangular entries",
                    "aᵢⱼ > 0: cooperation between species i and j;  aᵢⱼ < 0: competition",
                    "B is diagonal:  decay rate bᵢ for each species → 5 parameters",
                    "Total: 20 parameters organized in 5 blocks (M1–M5)",
                ],
                17,
            ),
        ],
    ),
    # ── 9. Biological Constraints (NEW) ────────────────────────
    (
        "Biological Constraints and Parameter Reduction",
        [
            ("h", "Experiment-based parameter locking (Nishioka Algorithm)"),
            (
                "t",
                "Lock 5 interaction parameters to zero where no interaction was observed (Heine et al., Fig. 4C):",
                18,
            ),
            ("g", 0.10),
            ("e", _EQ["locked_set"], 20),
            ("g", 0.10),
            (
                "b",
                [
                    "a₃₄ = 0:  Vei ↔ F.n — no direct metabolic pathway",
                    "a₂₃ = 0:  A.n ↔ Vei — no direct metabolic link",
                    "a₂₄ = 0:  A.n ↔ F.n — no direct interaction",
                    "a₁₅ = 0:  S.o ↔ P.g — no direct interaction",
                    "a₂₅ = 0:  A.n ↔ P.g — no direct interaction",
                ],
                16,
            ),
            ("g", 0.10),
            ("h", "Effective parameter dimension"),
            ("e", _EQ["n_free"], 24),
            ("g", 0.08),
            (
                "t",
                "→  Improved sampling efficiency + exclusion of biologically implausible estimates.",
                18,
            ),
        ],
    ),
    # ── 10. Hill gating ────────────────────────────────────────
    (
        "Hill-Function Gating Mechanism",
        [
            (
                "t",
                "F. nucleatum (species 4) acts as a bridging organism for P. gingivalis (species 5).",
                20,
            ),
            (
                "t",
                "Without sufficient F. nucleatum, interactions of P. gingivalis are suppressed.",
                18,
            ),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$x_{\mathrm{FN}} = \phi_4\,\psi_4$", "Effective biomass of F. nucleatum"),
                ],
            ),
            ("g", 0.10),
            ("tb", "Hill function"),
            ("g", 0.05),
            ("e", _EQ["hill"], 28),
            ("g", 0.08),
            (
                "s",
                [
                    (
                        r"$K_{\mathrm{hill}}$",
                        "Half-saturation constant — concentration at 50% gating",
                    ),
                    (r"$n_{\mathrm{hill}}$", "Hill coefficient — controls steepness of the switch"),
                ],
                18,
            ),
            ("g", 0.10),
            ("tb", "Effective interaction"),
            ("g", 0.05),
            ("e", _EQ["a_eff"], 24),
        ],
    ),
    # ── 11. Observation model ──────────────────────────────────
    (
        "Observation Model and Likelihood",
        [
            ("tb", "Observation model"),
            ("g", 0.08),
            ("e", _EQ["obs_model"], 24),
            ("g", 0.15),
            (
                "s",
                [
                    (r"$y_{k,s}$", "Observation at time t_k, species s"),
                    (r"$\hat{y}_{k,s}(\boldsymbol{\theta})$", "Model prediction via ROM"),
                    (r"$\sigma_{\mathrm{obs}}$", "Standard deviation of the observation noise"),
                ],
            ),
            ("g", 0.10),
            ("tb", "Log-likelihood"),
            ("g", 0.08),
            ("e", _EQ["loglik"], 24),
            ("g", 0.10),
            ("tb", "Prior distribution"),
            ("g", 0.08),
            ("e", _EQ["prior"], 22),
        ],
    ),
    # ── 12. ROM / TSM + Linearization ──────────────────────────
    (
        "ROM — TSM with Linearization Update",
        [
            ("tb", "Dimension reduction"),
            ("g", 0.08),
            ("e", _EQ["rom_approx"], 22),
            ("g", 0.08),
            (
                "t",
                "V is constructed from POD (Proper Orthogonal Decomposition) of FOM snapshots.",
                18,
            ),
            ("g", 0.12),
            ("tb", "Linearisation of FOM → time-stepping map"),
            ("g", 0.08),
            ("e", _EQ["rom_map"], 26),
            ("g", 0.12),
            ("h", "Linearization point update during TMCMC"),
            ("t", "At each TMCMC stage, update the linearization point θ₀ to the current MAP:", 17),
            ("g", 0.05),
            ("e", _EQ["linearization_update"], 20),
            ("g", 0.08),
            (
                "b",
                [
                    "ROM accuracy improves as the posterior concentrates",
                    "Skip update when ROM error is below threshold → reduces FOM evaluation cost",
                ],
                17,
            ),
        ],
    ),
    # ── 13. TMCMC Details ──────────────────────────────────────
    (
        "TMCMC — Algorithm Details",
        [
            ("tb", "Tempered posterior distribution"),
            ("g", 0.05),
            ("e", _EQ["tempered"], 20),
            ("g", 0.10),
            ("tc", "1)  Weight computation", ACCENT, 17, True),
            ("e", _EQ["weight_update"], 18),
            ("g", 0.05),
            ("tc", "2)  Adaptive β schedule", ACCENT, 17, True),
            ("t", "     Select β_{m+1} via bisection to maintain target ESS (N/2):", 16),
            ("e", _EQ["ess"], 18),
            ("g", 0.05),
            ("tc", "3)  Resampling + MCMC mutation", ACCENT, 17, True),
            ("t", "     Proposal distribution (Gelman optimal scaling 2.38²/d):", 16),
            ("e", _EQ["proposal_cov"], 18),
            ("g", 0.05),
            ("tc", "4)  MH acceptance probability", ACCENT, 17, True),
            ("e", _EQ["mh_accept"], 17),
        ],
    ),
    # ── 14. Sequential Estimation + Conditions (NEW) ───────────
    (
        "4-Stage Sequential Estimation and Experimental Conditions",
        [
            ("h", "4-stage sequential estimation"),
            (
                "b",
                [
                    "Stage 1 (M1):  S.o + A.n  → 5 params (a₁₁, a₁₂, a₂₂, b₁, b₂)",
                    "Stage 2 (M2):  Vei + F.n  → 5 params (a₃₃, a₃₄, a₄₄, b₃, b₄)",
                    "Stage 3 (M3+M4):  Cross + P.g self → 6 params (a₁₃, a₁₄, a₂₃, a₂₄, a₅₅, b₅)",
                    "Stage 4 (M5):  P.g cross → 4 params (a₁₅, a₂₅, a₃₅, a₄₅)",
                ],
                16,
            ),
            ("g", 0.05),
            ("t", "After each stage, fix MAP estimates in θ_base and proceed:", 17),
            ("e", _EQ["seq_update"], 20),
            ("g", 0.10),
            ("h", "4 experimental conditions (Heine et al.)"),
            (
                "b",
                [
                    "Commensal Static:  lock 9 / estimate 11 — pathogen-free baseline",
                    "Dysbiotic Static:   lock 5 / estimate 15 — pathogens present, static culture",
                    "Commensal HOBIC:  lock 8 / estimate 12 — flow culture, S.o growth estimated",
                    "Dysbiotic HOBIC:   lock 0 / estimate 20 — all unlocked (Surge reproduction)",
                ],
                16,
            ),
            ("g", 0.05),
            (
                "t",
                "→  Lock set varies per condition to ensure biologically valid estimation space.",
                17,
            ),
        ],
    ),
    # ── 15. Summary ────────────────────────────────────────────
    (
        "Summary and Outlook",
        [
            ("h", "Summary"),
            (
                "b",
                [
                    "FOM:  Nonlinear dissipative system for 5-species biofilm based on free energy.",
                    "Parameter reduction: biological constraints reduce 20 → 15 free parameters.",
                    "ROM (TSM):  Low-dim approximation + adaptive linearization point update.",
                    "TMCMC:  Adaptive β schedule, optimal proposal covariance, K-step mutation.",
                    "4-stage sequential estimation:  biologically-grouped stages for stable inference.",
                ],
                17,
            ),
            ("g", 0.15),
            ("h", "Outlook"),
            (
                "b",
                [
                    "Application to experimental data and validation of estimation accuracy.",
                    "Model evidence: Bayes factor comparison via TMCMC at no extra cost.",
                    "Uncertainty quantification:  credible intervals from the posterior.",
                    "Multi-condition estimation:  all 4 conditions (commensal/dysbiotic × static/HOBIC).",
                ],
                17,
            ),
            ("g", 0.10),
            ("h", "Model evidence (TMCMC byproduct)"),
            ("e", _EQ["evidence"], 20),
        ],
    ),
    # ── Backup: Hamilton → implementation ──────────────────────
    (
        "Backup: From Hamilton Principle to Implementation",
        [
            ("h", "Step 1:  Variational principle (Hamilton)"),
            ("e", _EQ["hamilton_var"], 24),
            ("g", 0.10),
            ("h", "Step 2:  Strong form (Euler-Lagrange)"),
            ("e", _EQ["dissipative"], 22),
            ("g", 0.10),
            ("h", "Step 3:  Time discretisation (implicit Euler)"),
            ("e", _EQ["implicit_euler_full"], 20),
            ("g", 0.10),
            ("h", "Step 4:  Residual equation + Newton iteration"),
            ("e", _EQ["residual_newton"], 18),
            ("g", 0.15),
            ("t", "→  Paper eqs. (16)-(18) correspond to these 4 steps.", 17),
            ("t", "→  The Q in Step 4 is the colour-coded equation from the earlier slide.", 17),
        ],
    ),
]


# ═══════════════════════════════════════════════════════════════
#  Build & save
# ═══════════════════════════════════════════════════════════════
def build(lang: str = "ja") -> Presentation:
    global _slide_counter
    _slide_counter = 0

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # title slide (custom layout)
    _slide_title(prs, lang)
    # title speaker note
    title_slide = prs.slides[0]
    notes = title_slide.notes_slide.notes_text_frame
    if lang == "ja":
        notes.text = (
            "本日は，5種バイオフィルムを対象に，自由エネルギーに基づくFOMと，"
            "ROM・TMCMCを組み合わせたパラメータ推定の枠組みをご紹介します。"
            "数理モデルの構造と推定アルゴリズムの設計に焦点を当てます。"
        )
    else:
        notes.text = (
            "This talk presents a free-energy-based FOM for a 5-species biofilm "
            "and a ROM+TMCMC framework for Bayesian parameter estimation. "
            "We focus on the model structure and algorithm design."
        )

    # content slides (data-driven)
    slides_data = SLIDES_JA if lang == "ja" else SLIDES_EN
    notes_data = NOTES_JA if lang == "ja" else NOTES_EN
    for idx, (title, elements) in enumerate(slides_data):
        slide = _new_slide(prs, title)
        _render_elements(slide, elements)
        if idx < len(notes_data) and notes_data[idx]:
            note_frame = slide.notes_slide.notes_text_frame
            note_frame.text = notes_data[idx]

    return prs


def main() -> None:
    for lang, suffix in [("ja", ""), ("en", "_en")]:
        print(f"Building {lang.upper()} presentation ...")
        prs = build(lang)
        path = os.path.join(SCRIPT_DIR, f"biofilm_fom_rom_tmcmc_nishioka{suffix}.pptx")
        prs.save(path)
        print(f"  Saved -> {path}")
    print("Done.")


if __name__ == "__main__":
    main()
