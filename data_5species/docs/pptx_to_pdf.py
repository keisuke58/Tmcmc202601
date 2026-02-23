#!/usr/bin/env python3
"""Convert PPTX to PDF preview using python-pptx + PIL + reportlab."""

import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas

PPTX = "/home/nishioka/IKM_Hiwi/Tmcmc202601/data_5species/docs/nishioka_biofilm_tmcmc_presentation.pptx"
PDF  = "/home/nishioka/IKM_Hiwi/Tmcmc202601/data_5species/docs/nishioka_biofilm_tmcmc_presentation.pdf"

# Render scale: PPTX is 13.333 x 7.5 inches -> render at 150 DPI
DPI = 150
prs = Presentation(PPTX)
sw_emu = prs.slide_width
sh_emu = prs.slide_height
sw_px = int(sw_emu / 914400 * DPI)
sh_px = int(sh_emu / 914400 * DPI)

print(f"Slide size: {sw_emu/914400:.1f} x {sh_emu/914400:.1f} inches")
print(f"Render size: {sw_px} x {sh_px} px at {DPI} DPI")
print(f"Total slides: {len(prs.slides)}")

def emu_to_px(emu):
    return int(emu / 914400 * DPI)

def rgb_from_pptx(color):
    """Extract RGB tuple from pptx color."""
    try:
        if color and color.rgb:
            r = color.rgb.__str__()
            return (int(r[0:2], 16), int(r[2:4], 16), int(r[4:6], 16))
    except:
        pass
    return None

def get_fill_color(shape):
    """Get fill color of a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                r = fc.rgb.__str__()
                return (int(r[0:2], 16), int(r[2:4], 16), int(r[4:6], 16))
    except:
        pass
    return None

def get_bg_color(slide):
    """Get slide background color."""
    try:
        bg = slide.background.fill
        if bg.type is not None:
            fc = bg.fore_color
            if fc and fc.rgb:
                r = fc.rgb.__str__()
                return (int(r[0:2], 16), int(r[2:4], 16), int(r[4:6], 16))
    except:
        pass
    return (240, 242, 245)  # default SLIDE_BG

# Try to find a usable font
def find_font(size=14, bold=False):
    font_paths = [
        "/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/google-noto/NotoSans-Regular.ttf",
        "/usr/share/fonts/liberation-sans/LiberationSans-Regular.ttf",
        "/usr/share/fonts/dejavu-sans/DejaVuSans.ttf",
        "/usr/share/fonts/gnu-free/FreeSans.ttf",
    ]
    bold_paths = [
        "/usr/share/fonts/google-noto-cjk/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/google-noto/NotoSans-Bold.ttf",
        "/usr/share/fonts/liberation-sans/LiberationSans-Bold.ttf",
        "/usr/share/fonts/dejavu-sans/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/gnu-free/FreeSansBold.ttf",
    ]
    paths = bold_paths if bold else font_paths
    paths += font_paths  # fallback to regular if bold not found
    for fp in paths:
        if os.path.exists(fp):
            try:
                return ImageFont.truetype(fp, size)
            except:
                continue
    return ImageFont.load_default()

# Pre-cache some fonts
font_cache = {}
def get_font(size, bold=False):
    key = (size, bold)
    if key not in font_cache:
        font_cache[key] = find_font(size, bold)
    return font_cache[key]

slide_images = []

for slide_idx, slide in enumerate(prs.slides):
    bg = get_bg_color(slide)
    img = Image.new("RGB", (sw_px, sh_px), bg)
    draw = ImageDraw.Draw(img)

    # Sort shapes by z-order (approximate: just iterate in order)
    shapes = list(slide.shapes)

    for shape in shapes:
        x = emu_to_px(shape.left) if shape.left else 0
        y = emu_to_px(shape.top) if shape.top else 0
        w = emu_to_px(shape.width) if shape.width else 0
        h = emu_to_px(shape.height) if shape.height else 0

        # Draw shape fill (rectangles, rounded rects)
        if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            fc = get_fill_color(shape)
            if fc:
                draw.rectangle([x, y, x + w, y + h], fill=fc)
            # Draw border
            try:
                line = shape.line
                if line and line.color and line.color.rgb:
                    lc = line.color.rgb.__str__()
                    border = (int(lc[0:2], 16), int(lc[2:4], 16), int(lc[4:6], 16))
                    lw = max(1, int((line.width or Pt(1)) / 914400 * DPI))
                    draw.rectangle([x, y, x + w, y + h], outline=border, width=lw)
            except:
                pass

        # Draw pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_blob = shape.image.blob
                pic = Image.open(io.BytesIO(img_blob))
                pic = pic.convert("RGBA")
                pic_resized = pic.resize((max(1, w), max(1, h)), Image.LANCZOS)
                img.paste(pic_resized, (x, y), pic_resized if pic_resized.mode == "RGBA" else None)
            except Exception as e:
                draw.rectangle([x, y, x + w, y + h], outline=(200, 200, 200))
                draw.text((x + 5, y + 5), f"[img err]", fill=(200, 0, 0),
                          font=get_font(10))

        # Draw text
        if shape.has_text_frame:
            tf = shape.text_frame
            ty = y
            for para in tf.paragraphs:
                text = para.text
                if not text.strip():
                    ty += 4
                    continue

                # Get font properties from first run
                sz = 14
                bold = False
                color = (45, 45, 45)
                try:
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            sz = int(run.font.size / 12700)  # EMU to pt approx
                        bold = run.font.bold or False
                        fc = rgb_from_pptx(run.font.color)
                        if fc:
                            color = fc
                    elif para.font and para.font.size:
                        sz = int(para.font.size / 12700)
                except:
                    pass

                font_px = max(8, int(sz * DPI / 72))
                font = get_font(font_px, bold)

                # Simple text wrapping
                max_w = w - 4
                lines = []
                for raw_line in text.split('\n'):
                    if not raw_line:
                        lines.append("")
                        continue
                    words = raw_line.split(' ')
                    current = ""
                    for word in words:
                        test = f"{current} {word}".strip()
                        bbox = font.getbbox(test)
                        tw = bbox[2] - bbox[0] if bbox else 0
                        if tw > max_w and current:
                            lines.append(current)
                            current = word
                        else:
                            current = test
                    if current:
                        lines.append(current)

                for line in lines:
                    if ty + font_px > y + h:
                        break
                    # Handle alignment
                    tx = x + 2
                    try:
                        from pptx.enum.text import PP_ALIGN
                        if para.alignment == PP_ALIGN.CENTER:
                            bbox = font.getbbox(line)
                            tw = bbox[2] - bbox[0] if bbox else 0
                            tx = x + (w - tw) // 2
                        elif para.alignment == PP_ALIGN.RIGHT:
                            bbox = font.getbbox(line)
                            tw = bbox[2] - bbox[0] if bbox else 0
                            tx = x + w - tw - 4
                    except:
                        pass
                    draw.text((tx, ty), line, fill=color, font=font)
                    ty += font_px + 2

    slide_images.append(img)
    print(f"  Rendered slide {slide_idx + 1}/{len(prs.slides)}")

# Create PDF
print(f"\nCreating PDF...")
page_w = sw_px
page_h = sh_px
c = canvas.Canvas(PDF, pagesize=(page_w, page_h))

for i, img in enumerate(slide_images):
    # Save slide image to bytes
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    from reportlab.lib.utils import ImageReader
    ir = ImageReader(buf)
    c.drawImage(ir, 0, 0, width=page_w, height=page_h)
    c.showPage()

c.save()
print(f"PDF saved: {PDF}")
print(f"Total pages: {len(slide_images)}")
